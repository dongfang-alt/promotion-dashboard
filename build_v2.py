#!/usr/bin/env python3
"""
2026推广年中会议汇报 — 基于XMind新框架重写
专家视角：重构逻辑线、升级结论、精简表达
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ═══ 设计系统 ═══
BG      = RGBColor(0x0A,0x0F,0x1E)
CARD    = RGBColor(0x14,0x1E,0x33)
HEADER  = RGBColor(0x0F,0x17,0x2A)
ACCENT  = RGBColor(0x1E,0x29,0x3B)

BLUE    = RGBColor(0x3B,0x82,0xF6); BLUE_BG = RGBColor(0x1E,0x3A,0x5F)
GREEN   = RGBColor(0x22,0xC5,0x5E); GREEN_BG= RGBColor(0x14,0x52,0x2D)
AMBER   = RGBColor(0xF5,0x9E,0x0B); AMBER_BG= RGBColor(0x78,0x3E,0x0F)
RED     = RGBColor(0xEF,0x44,0x44); RED_BG  = RGBColor(0x7F,0x1D,0x1D)
PURPLE  = RGBColor(0x8B,0x5C,0xF6); PURPLE_BG=RGBColor(0x3B,0x1F,0x6E)
TEAL    = RGBColor(0x14,0xB8,0xA6)

WHITE   = RGBColor(0xF8,0xFA,0xFC)
GREY    = RGBColor(0x94,0xA3,0xB8)
GREY2   = RGBColor(0x64,0x74,0x8B)
BORDER  = RGBColor(0x33,0x40,0x55)

SW = 12192000; SH = 6858000; M = 720000; G = 180000

prs = Presentation()
prs.slide_width = SW; prs.slide_height = SH
BL = prs.slide_layouts[6]

# ═══ 工具 ═══
def bg(slide, c=BG):
    s=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,SW,SH)
    s.fill.solid(); s.fill.fore_color.rgb=c; s.line.fill.background()

def rect(slide,l,t,w,h,fill=None,border=None,bw=None,r=None):
    s=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if r else MSO_SHAPE.RECTANGLE,l,t,w,h)
    if fill: s.fill.solid(); s.fill.fore_color.rgb=fill
    else: s.fill.background()
    if border: s.line.color.rgb=border
    if bw: s.line.width=bw
    else: s.line.fill.background()
    return s

def tb(slide,l,t,w,h,text,fs=14,c=WHITE,b=False,al=PP_ALIGN.LEFT,fn='微软雅黑',ls=1.2):
    tx=slide.shapes.add_textbox(l,t,w,h)
    tf=tx.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]; p.text=text
    p.font.size=Pt(fs); p.font.color.rgb=c; p.font.bold=b; p.font.name=fn; p.alignment=al
    p.space_after=Pt(1); p.line_spacing=Pt(fs*ls)
    return tx

def mtb(slide,l,t,w,h,lines,fs=13,c=GREY,fn='微软雅黑',ls=1.3):
    tx=slide.shapes.add_textbox(l,t,w,h)
    tf=tx.text_frame; tf.word_wrap=True
    for i,(text,opts) in enumerate(lines):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        sz=opts.get('size',fs); cl=opts.get('color',c); bo=opts.get('bold',False)
        p.text=text; p.font.size=Pt(sz); p.font.color.rgb=cl; p.font.bold=bo
        p.font.name=fn; p.space_after=Pt(2); p.line_spacing=Pt(sz*ls)
    return tx

def hdr(slide,title,sub=""):
    rect(slide,0,0,SW,550000,fill=HEADER)
    rect(slide,0,550000,SW,20000,fill=BLUE)
    tb(slide,M,110000,7000000,330000,title,fs=24,c=WHITE,b=True)
    if sub: tb(slide,SW-M-3800000,140000,3800000,250000,sub,fs=11,c=GREY,al=PP_ALIGN.RIGHT)

def kpi(slide,l,t,w,label,value,vc=None,sub="",bg=CARD,ls=9,vs=14):
    rect(slide,l,t,w,480000,fill=bg,r=60000)
    tb(slide,l+70000,t+50000,w-140000,130000,label,fs=ls,c=GREY)
    tb(slide,l+70000,t+170000,w-140000,190000,value,fs=vs,c=vc or WHITE,b=True)
    if sub: tb(slide,l+70000,t+360000,w-140000,100000,sub,fs=8,c=GREY)

def div(slide,l,t,w): rect(slide,l,t,w,10000,fill=BORDER)

def sect(slide,l,t,w,num,title,color=BLUE):
    rect(slide,l,t,40000,250000,fill=color)
    tb(slide,l+80000,t+5000,500000,240000,num,fs=20,c=color,b=True)
    tb(slide,l+500000,t+5000,w-500000,240000,title,fs=17,c=WHITE,b=True)

# ═══════════════════════════════════
# SLIDE 1: 封面
# ═══════════════════════════════════
s=prs.slides.add_slide(BL); bg(s)
rect(s,0,0,SW,40000,fill=BLUE)
rect(s,M,1200000,35000,2400000,fill=BLUE)
tb(s,M+160000,1200000,9000000,720000,'2026 推广年中会议汇报',fs=44,c=WHITE,b=True)
tb(s,M+160000,1920000,8000000,300000,'E-commerce Promotion Mid-Year Review',fs=16,c=GREY)
tb(s,M+160000,2320000,7000000,320000,'H1 平台复盘  /  H2 三大规划  /  组织升级',fs=20,c=BLUE,b=True)
div(s,M+160000,2740000,6200000)
y0=4800000
tb(s,M+160000,y0,3000000,260000,'汇报人：东方',fs=13,c=GREY)
tb(s,M+160000,y0+300000,3000000,260000,'部门：国内营销中心',fs=13,c=GREY)
tb(s,M+160000,y0+600000,3000000,260000,'汇报时间：2026.07.23',fs=13,c=GREY)
tb(s,SW-M-3800000,SH-M-300000,3800000,300000,'从投放执行到增长经营',fs=15,c=BLUE,b=True,al=PP_ALIGN.RIGHT)

# ═══════════════════════════════════
# SLIDE 2: H1数据复盘 — 天猫+京东核心指标
# ═══════════════════════════════════
s=prs.slides.add_slide(BL); bg(s)
hdr(s,'上半年复盘：核心店铺数据','2026H1 vs 2025H1（同比）/ 2026H1 vs 2025H2（环比）')

tb(s,M,650000,10600000,250000,'重点分析店铺：天猫旗舰店、京东自营  |  ROI用数值，其他用百分比',fs=11,c=GREY)

# ── 天猫旗舰店行 ──
y=950000; lw=1550000; kw=1050000
rect(s,M,y,lw,1020000,fill=BLUE_BG,r=80000)
tb(s,M+100000,y+130000,lw-200000,240000,'天猫旗舰店',fs=15,c=WHITE,b=True)
tb(s,M+100000,y+420000,lw-200000,180000,'重点店铺',fs=10,c=BLUE,b=True)
tb(s,M+100000,y+650000,lw-200000,280000,'H1：以费换量，ROI修复明显',fs=9,c=GREY)

tm_kpi=[("消耗同比","+30.3%",AMBER),("消耗环比","-0.8%",GREY),
        ("ROI同比","+39.3%",GREEN),("ROI环比","+10.4%",GREEN),
        ("费比同比","+2.6pt",AMBER),("费比环比","-0.9pt",GREEN),
        ("销售同比","+10.0%",GREEN),("销售环比","+4.4%",GREEN)]
for i,(lb,vl,vc) in enumerate(tm_kpi):
    x=M+lw+70000+i*(kw+35000)
    kpi(s,x,y,kw,lb,vl,vc)

# ── 京东自营行 ──
y2=2070000
rect(s,M,y2,lw,1020000,fill=AMBER_BG,r=80000)
tb(s,M+100000,y2+130000,lw-200000,240000,'京东自营',fs=15,c=WHITE,b=True)
tb(s,M+100000,y2+420000,lw-200000,180000,'重点店铺',fs=10,c=AMBER,b=True)
tb(s,M+100000,y2+650000,lw-200000,280000,'H1：控费显著，规模承压',fs=9,c=GREY)

jd_kpi=[("消耗同比","-13.2%",GREY),("消耗环比","-15.8%",GREY),
        ("ROI同比","+25.3%",GREEN),("ROI环比","+36.4%",GREEN),
        ("费比同比","+0.3pt",AMBER),("费比环比","-1.0pt",GREEN),
        ("销售同比","-15.8%",RED),("销售环比","-8.1%",RED)]
for i,(lb,vl,vc) in enumerate(jd_kpi):
    x=M+lw+70000+i*(kw+35000)
    kpi(s,x,y2,kw,lb,vl,vc)

# ── 核心矛盾卡片 ──
y3=3300000; cw=5000000
# 天猫总结
rect(s,M,y3,cw,1100000,fill=CARD,r=100000)
sect(s,M+150000,y3+80000,cw-300000,'','天猫旗舰店：以费换量，ROI修复成立',BLUE)
mtb(s,M+150000,y3+380000,cw-300000,680000,[
    ('▸ ROI同比+39.3%、环比+10.4%，推广优化效果明确；消耗同比+30.3%，平台联合放量明显',{}),
    ('▸ Q2联合平台签订高费比合同，ROI/CPC下滑可控的前提下，销售同比+9.6%、环比+24.4%',{}),
    ('▸ H2重点：保爆款、控节奏、抢核心词，站外引流必须回看站内承接',{}),
    ('▸ 同期新开设：天猫数码旗舰店（新起盘，轻预算测试）',{}),
],fs=11,c=GREY)

# 京东总结
rect(s,M+cw+200000,y3,cw+300000,1100000,fill=CARD,r=100000)
sect(s,M+cw+350000,y3+80000,cw,'','京东自营：控费显著，但规模弹性丧失',AMBER)
mtb(s,M+cw+350000,y3+380000,cw,680000,[
    ('▸ ROI同比+25.3%、环比+36.4%，H1推广优化效果非常明确',{}),
    ('▸ 但销售同比-15.8%、环比-8.1%，说明控费已到临界点，不能只节流不开源',{}),
    ('▸ H2重点：在高ROI计划中恢复核心品类扩量，关键词质量分、人群分层、爆品预算保护',{}),
    ('▸ 京东POP、Mr.Green自营作为结构补充，ROI修复同样明确',{}),
],fs=11,c=GREY)

# 底部结论
y4=4650000
rect(s,M,y4,SW-M*2,480000,fill=ACCENT,r=80000)
tb(s,M+200000,y4+50000,SW-M*2-400000,180000,'H1 核心判断',fs=14,c=BLUE,b=True)
tb(s,M+200000,y4+240000,SW-M*2-400000,200000,'H1完成了推广效率的全面修复（ROI↑、费比↓），但京东销售弹性已丧失，H2必须从「止血」转向「造血」——在天猫用可控费比换确定性增长，在京东恢复效率与规模的平衡。',fs=11,c=GREY)

# ═══════════════════════════════════
# SLIDE 3: H1体系成果 + 压力点
# ═══════════════════════════════════
s=prs.slides.add_slide(BL); bg(s)
hdr(s,'上半年复盘：体系成果与压力点','H1 除了数据之外，我们还建立了什么？还缺什么？')

# ── 体系成果（左） ──
y=750000
tb(s,M+200000,y,5000000,350000,'体系成果：从个人经验到团队能力',fs=18,c=GREEN,b=True)

achievements=[
    ('推广系统培训','指导运营提升理论与实操能力','让运营理解ROI/CPC/CVR联动关系，从只会建计划到会看趋势',GREEN),
    ('多维数据监控表','制定复盘模板，周期复盘、及时调优','异常发现周期从「月度汇总」压缩到「周度扫描」，问题响应提速',BLUE),
    ('推广新人培养','部门助理铁蛋，半年培养为推广助理并成功转岗','验证了方法论可复制，证明了培养体系的有效性',PURPLE),
]
for i,(t,sub,desc,cl) in enumerate(achievements):
    ay=y+450000+i*750000
    rect(s,M,ay,SW-M*2,650000,fill=CARD,r=80000)
    rect(s,M,ay,40000,650000,fill=cl)
    tb(s,M+130000,ay+60000,3200000,220000,t,fs=14,c=cl,b=True)
    tb(s,M+130000,ay+280000,SW-M*2-260000,180000,sub,fs=11,c=WHITE)
    tb(s,M+130000,ay+460000,SW-M*2-260000,160000,desc,fs=10,c=GREY)

# ── 压力点（右下方，叠在体系成果下面）──
py=2950000
tb(s,M+200000,py,5000000,320000,'当前最大压力点：制约提效的三个瓶颈',fs=18,c=AMBER,b=True)

pains=[
    ('长尾品优化难','虽设定了商品等级、抓重点商品，但1人统筹多店铺多模式，精力有限，长尾品和细分计划难以持续跟踪优化',AMBER),
    ('数据分析工具缺失','推广商品数据量大、维度多，缺少系统化、多维度的数据处理工具，人工拉表效率低、易遗漏',RED),
    ('全链路断层','从发现问题→归因→策略→执行→复盘，缺少专业深度的全链路分析指导，决策依赖个人经验',RED),
]
for i,(t,desc,cl) in enumerate(pains):
    pay=py+400000+i*600000
    rect(s,M,pay,SW-M*2,520000,fill=CARD,r=80000)
    rect(s,M,pay,40000,520000,fill=cl)
    tb(s,M+130000,pay+50000,3500000,200000,t,fs=13,c=cl,b=True)
    tb(s,M+130000,pay+260000,SW-M*2-260000,220000,desc,fs=10,c=GREY)

# 底部承接
by=5250000
rect(s,M,by,SW-M*2,380000,fill=ACCENT,r=80000)
tb(s,M+200000,by+40000,SW-M*2-400000,300000,'这三个压力点，恰好是 H2 规划的方向：扩大监控范围 → 引入 AI 工具 → 建立标准化分析体系',fs=12,c=WHITE,b=True)

# ═══════════════════════════════════
# SLIDE 4: H2-01 推广精细化提效
# ═══════════════════════════════════
s=prs.slides.add_slide(BL); bg(s)
hdr(s,'下半年规划 01：推广精细化提效','从核心品类关注到全维度监控调优')

tb(s,M,650000,10600000,240000,'背景：H1仅1人统筹多店铺、多模式，长尾品难覆盖；H2通过推广助理+AI工具，将监控范围从核心品类扩展到全品类、全计划。',fs=11,c=GREY)

# 四层级商品策略
y=1000000; qw=2620000; qh=1800000
quads=[
    ('核心爆款',BLUE,BLUE_BG,'保排名·保转化·保预算','看 ROI环比、核心词排名、活动承接'),
    ('潜力单品',GREEN,GREEN_BG,'测人群·测素材·测场景','看 CPC、CVR、加购率、ROI趋势'),
    ('新品',AMBER,AMBER_BG,'小预算冷启·先跑模型','看 点击率、加购率、素材胜出率'),
    ('中长尾品',PURPLE,PURPLE_BG,'低预算验证·快速止损','看 ROI阈值、转化率、费用泄露'),
]
for i,(t,cl,bg_cl,act,watch) in enumerate(quads):
    x=M+i*(qw+170000)
    rect(s,x,y,qw,qh,fill=bg_cl,r=100000)
    tb(s,x+150000,y+120000,qw-300000,260000,t,fs=17,c=cl,b=True)
    tb(s,x+150000,y+420000,qw-300000,220000,act,fs=12,c=WHITE,b=True)
    tb(s,x+150000,y+680000,qw-300000,450000,watch,fs=10,c=GREY)
    rect(s,x,y,qw,40000,fill=cl)

# 监控升级三列
y2=3000000; c3w=3400000
upgrades=[
    ('1. 监控范围升级',BLUE,BLUE_BG,
     ['过去：核心品类、爆款商品、主要账户优先', '升级：平台→店铺→品类→商品→计划→关键词→人群→素材','目标：从「事后发现大问题」→「过程识别小异常」']),
    ('2. 人力分工升级',GREEN,GREEN_BG,
     ['推广负责人：策略、预算、平台打法、关键决策','推广助理：数据整理、异常初筛、素材/计划跟进','目标：负责人从拉表中释放，投入策略与复盘']),
    ('3. AI工具提效',PURPLE,PURPLE_BG,
     ['AI：多维数据处理、异常识别、趋势归因、策略草案','人工：业务判断、动作取舍、执行落地、复盘纠偏','目标：提速诊断，降低对个人经验的依赖']),
]
c3s=(SW-c3w*3-200000*2)//2
for i,(t,cl,bg_cl,items) in enumerate(upgrades):
    x=c3s+i*(c3w+200000)
    rect(s,x,y2,c3w,1850000,fill=bg_cl,r=120000)
    rect(s,x,y2,c3w,40000,fill=cl)
    tb(s,x+180000,y2+120000,c3w-360000,260000,t,fs=15,c=cl,b=True)
    ils=[(f'• {it}',{'size':10,'color':GREY}) for it in items]
    mtb(s,x+180000,y2+450000,c3w-360000,1300000,ils,fs=10,c=GREY)

# 落地机制
by=5100000
rect(s,M,by,SW-M*2,380000,fill=ACCENT,r=80000)
tb(s,M+200000,by+35000,SW-M*2-400000,300000,'落地机制：日看异常 → 周看商品层级迁移 → 月看平台结构。助理负责数据+跟进，AI负责初筛+策略草案，负责人负责判断+关键动作。',fs=11,c=WHITE,b=True)

# ═══════════════════════════════════
# SLIDE 5: H2-02 赋能品类增长
# ═══════════════════════════════════
s=prs.slides.add_slide(BL); bg(s)
hdr(s,'下半年规划 02：赋能品类增长','推广组合打法+数据反馈，赋能品类销售增长')

tb(s,M,650000,10600000,280000,'目标：推广不只是买流量——通过投放组合和数据反馈，帮助品类找到确定的增长路径，指导货品、价格、内容和活动节奏。',fs=12,c=GREY)

# 三列大卡
y=1050000; bcw=3400000; bch=3800000
bcs=(SW-bcw*3-200000*2)//2
bigs=[
    ('数据反馈','识别机会',GREEN,
     ['按品类/商品监控 ROI、CPC、CVR、点击率、收藏加购趋势',
      '识别三类机会：低CPC高CVR→放量；高点击低转化→承接问题；高消耗低ROI→止损',
      '输出给运营：哪些品类加预算，哪些商品需优化页面/价格/权益']),
    ('组合打法','提升效率',BLUE,
     ['搜索承接明确需求，推荐/场景拓新人群，站外种草补认知',
      '核心品→搜索+推荐联动；潜力品→小预算多场景测试；新品→先内容种草再站内承接',
      '每类打法绑定复盘指标，不看动作看结果']),
    ('经营协同','放大结果',PURPLE,
     ['推广结论同步商品/运营/供应链：爆款备货、价格策略、详情页承接、活动资源',
      '大促前货品池+预算池联审，避免有流量无库存/权益/承接',
      '沉淀品类打法，形成可复用增长模型']),
]
for i,(t,sub,cl,items) in enumerate(bigs):
    x=bcs+i*(bcw+200000)
    rect(s,x,y,bcw,bch,fill=CARD,r=120000)
    rect(s,x,y,bcw,45000,fill=cl)
    tb(s,x+220000,y+130000,bcw-440000,280000,t,fs=19,c=cl,b=True)
    tb(s,x+220000,y+400000,bcw-440000,240000,sub,fs=14,c=WHITE,b=True)
    ils=[(f'{j+1}. {it}',{'size':11,'color':GREY}) for j,it in enumerate(items)]
    mtb(s,x+220000,y+750000,bcw-440000,2900000,ils,fs=11,c=GREY,ls=1.4)

# 底部
by=5150000
rect(s,M,by,SW-M*2,380000,fill=ACCENT,r=80000)
tb(s,M+200000,by+40000,SW-M*2-400000,300000,'核心输出：从「投放复盘表」升级为「品类增长建议表」，让推广数据真正指导经营动作。',fs=13,c=WHITE,b=True)

# ═══════════════════════════════════
# SLIDE 6: H2-03 渠道策略
# ═══════════════════════════════════
s=prs.slides.add_slide(BL); bg(s)
hdr(s,'下半年规划 03：渠道策略方向','天猫 / 京东 / 拼多多')

tb(s,M,650000,10600000,260000,'核心逻辑：不再用同一套ROI标准一刀切，按平台角色分配目标、预算和考核指标。',fs=11,c=GREY)

y=1050000; chw=3400000; chh=4300000
chs=(SW-chw*3-200000*2)//2
chs_data=[
    ('天猫',BLUE,BLUE_BG,'可控费比换确定性增长',
     ['角色：品牌主阵地，市场份额承接渠道',
      '策略：围绕平台联合资源+活动节点+核心爆款，阶段性接受较高费比换取曝光与销售增长',
      '动作：核心词卡位、推荐扩量、UD站外引流、内容种草回流',
      '边界：高费比必须绑定活动目标+商品池+承接结果，避免无目标放量',
      'Q2验证：联合平台高费比合同，销售增长明显而ROI下滑可控']),
    ('京东',AMBER,AMBER_BG,'从过度控费到效率规模平衡',
     ['角色：效率修复与利润稳定渠道',
      '策略：守住高ROI计划，同时恢复核心品类与爆品的扩量弹性',
      '动作：关键词质量分、人群分层、低效词清理、爆品预算保护',
      '边界：优先看ROI环比、费比环比、CPC/CVR变化，防止只控费不增长',
      '核心矛盾：H1销售-16%，控费已到临界点']),
    ('拼多多',PURPLE,PURPLE_BG,'最小成本跑通方法论',
     ['角色：9月接手后的新增提效场',
      '策略：不急于规模放量，先完成账户体检、商品分层、计划结构搭建',
      '动作：搜索/场景基础计划、可投商品池、素材测试、ROI阈值+止损机制',
      '边界：先跑通方法，再放大预算，避免接手初期粗放消耗',
      '节奏：9月诊断→10月测试→双11小规模验证→12月沉淀规则']),
]
for i,(name,cl,bg_cl,sub,items) in enumerate(chs_data):
    x=chs+i*(chw+200000)
    rect(s,x,y,chw,chh,fill=bg_cl,r=120000)
    rect(s,x,y,chw,50000,fill=cl)
    tb(s,x+200000,y+120000,chw-400000,280000,name,fs=22,c=cl,b=True)
    tb(s,x+200000,y+400000,chw-400000,220000,sub,fs=13,c=WHITE,b=True)
    ils=[(f'▸ {it}',{'size':10,'color':GREY}) for it in items]
    mtb(s,x+200000,y+750000,chw-400000,3400000,ils,fs=10,c=GREY,ls=1.3)

# 底部
by=5600000
rect(s,M,by,SW-M*2,380000,fill=ACCENT,r=80000)
tb(s,M+200000,by+40000,SW-M*2-400000,300000,'渠道判断：天猫用可控费比换增长（看承接），京东从控费恢复到效率与规模平衡（补弹性），拼多多先建规则再规模化（控风险）。',fs=12,c=WHITE,b=True)

# ═══════════════════════════════════
# SLIDE 7: 组织升级01 — 成员管理
# ═══════════════════════════════════
s=prs.slides.add_slide(BL); bg(s)
hdr(s,'组织升级 01：组织成员管理','能力建设 / 分工机制 / 降低关键人员依赖')

tb(s,M,650000,10600000,260000,'核心目标：把个人经验拆成团队可执行的方法，降低业务对单一关键人员行业经验的依赖。',fs=11,c=GREY)

y=1050000; ocw=3400000; och=3800000
ocs=(SW-ocw*3-200000*2)//2
orgs=[
    ('能力短板→成长方向',BLUE,
     ['当前：成员能做基础执行，但对平台机制、商品层级、预算节奏的判断不足',
      '方向：从「会建计划、会调价」升级为「会看趋势、会归因、会判断扩量/止损」',
      '重点：ROI/费比/CPC/CVR联动分析，商品池分层，活动与推广节奏匹配']),
    ('负责人+助理协同',GREEN,
     ['负责人：平台策略、预算分配、关键账户诊断、跨部门协同',
      '助理：日报数据、异常初筛、基础计划跟进、素材与关键词维护',
      '协同：助理先筛异常→负责人判断优先级给动作→复盘沉淀为SOP']),
    ('标准沉淀：减少经验黑箱',PURPLE,
     ['商品分层SOP：核心爆款/潜力单品/新品/中长尾品→对应投放规则',
      '异常处理清单：ROI下滑/费比上升/CPC异常/CVR波动/消耗突增→排查路径',
      '复盘机制：周看计划、月看平台、活动看商品与承接→经验转团队资产']),
]
for i,(t,cl,items) in enumerate(orgs):
    x=ocs+i*(ocw+200000)
    rect(s,x,y,ocw,och,fill=CARD,r=120000)
    rect(s,x,y,ocw,45000,fill=cl)
    tb(s,x+200000,y+130000,ocw-400000,280000,t,fs=16,c=cl,b=True)
    ils=[(f'• {it}',{'size':10,'color':GREY}) for it in items]
    mtb(s,x+200000,y+500000,ocw-400000,3100000,ils,fs=10,c=GREY,ls=1.3)

by=5150000
rect(s,M,by,SW-M*2,380000,fill=ACCENT,r=80000)
tb(s,M+200000,by+40000,SW-M*2-400000,300000,'组织升级的本质：不是增加人手，是用清晰分工+标准流程，把团队从「被动执行」推向「主动诊断、主动优化」。',fs=12,c=WHITE,b=True)

# ═══════════════════════════════════
# SLIDE 8: 组织升级02 — AI赋能
# ═══════════════════════════════════
s=prs.slides.add_slide(BL); bg(s)
hdr(s,'组织升级 02：AI 赋能团队','AI策略与实际工作深度结合')

tb(s,M,650000,10600000,280000,'目标：AI不替代推广判断，而是提升数据处理、问题定位和策略草案效率——AI出数据+策略草案，人工做业务判断+执行落地+复盘纠偏。',fs=12,c=GREY)

y=1100000
ais=[
    ('AI多维数据处理',PURPLE,
     ['结构化平台/店铺/商品/计划/关键词数据，自动识别消耗、ROI、CPC、CVR异常',
      '从「人工拉表找问题」→「AI先筛异常，人工判断原因」',
      '优先落地：日报异常提醒 + 周度复盘摘要']),
    ('AI快速定位核心问题',BLUE,
     ['由总到细：平台→店铺→商品→计划→关键词，逐层下钻',
      '建立核心指标模块：消耗/ROI/CPC/CVR/点击率/加购/转化',
      'AI输出可能原因：流量成本、转化承接、预算节奏、商品竞争力、活动资源']),
    ('AI专业策略输出',GREEN,
     ['沉淀推广指令库：不同异常场景→对应分析指令+优化建议',
      '运营提组合问题→AI生成策略草案→负责人校验→指导实操',
      '降低推广指导成本，让团队更快形成标准化判断能力']),
]
for i,(t,cl,items) in enumerate(ais):
    x=ocs+i*(ocw+200000)
    rect(s,x,y,ocw,3200000,fill=CARD,r=120000)
    rect(s,x,y,ocw,45000,fill=cl)
    tb(s,x+200000,y+130000,ocw-400000,280000,t,fs=15,c=cl,b=True)
    ils=[(f'▸ {it}',{'size':10,'color':GREY}) for it in items]
    mtb(s,x+200000,y+500000,ocw-400000,2500000,ils,fs=10,c=GREY,ls=1.3)

# AI闭环
by=4600000
rect(s,M,by,SW-M*2,420000,fill=ACCENT,r=80000)
tb(s,M+200000,by+50000,SW-M*2-400000,300000,'AI闭环：数据清洗 → 异常识别 → 策略草案 → 人工判断 → 执行复盘 → 指令库沉淀',fs=13,c=WHITE,b=True)

# ═══════════════════════════════════
# SLIDE 9: H2行动节奏
# ═══════════════════════════════════
s=prs.slides.add_slide(BL); bg(s)
hdr(s,'H2 行动节奏','MILESTONE')

tb(s,M,650000,10600000,240000,'先建规则 → 接渠道 → 打大促 → 沉淀方法',fs=12,c=GREY)

y=1050000; tlw=2550000; tlh=1800000
tls=(SW-tlw*4-180000*3)//2
tls_data=[
    ('7-8月',BLUE,['完善商品分层规则、周度复盘模板','明确天猫/京东预算阈值+重点商品池','推广助理+AI工具落地，跑通日/周监控流程']),
    ('9月',PURPLE,['接手拼多多，账户诊断+商品分层+计划结构','搜索/场景基础投放规则搭建','双11预算框架+货品池初筛']),
    ('10-11月',AMBER,['双11预算节奏执行：爆品保护、低效止损','大促后复盘：平台/商品/素材/承接','输出大促推广复盘报告']),
    ('12月',GREEN,['沉淀年度推广手册+AI指令库+品类打法库','输出2027年预算与渠道规划建议','团队年度能力评估与下一年培养计划']),
]
for i,(mo,cl,items) in enumerate(tls_data):
    x=tls+i*(tlw+180000)
    rect(s,x,y,tlw,tlh,fill=CARD,r=120000)
    rect(s,x,y,tlw,50000,fill=cl)
    tb(s,x+180000,y+140000,tlw-360000,300000,mo,fs=20,c=cl,b=True)
    ils=[(f'• {it}',{'size':10,'color':GREY}) for it in items]
    mtb(s,x+180000,y+550000,tlw-360000,1100000,ils,fs=10,c=GREY,ls=1.3)

# 箭头
for i in range(3):
    ax=tls+(i+1)*tlw+i*180000-160000
    tb(s,ax,1800000,160000,300000,'→',fs=20,c=GREY2,b=True,al=PP_ALIGN.CENTER)

# ═══════════════════════════════════
# SLIDE 10: 团队协同 + 补充
# ═══════════════════════════════════
s=prs.slides.add_slide(BL); bg(s)
hdr(s,'团队协同机制与补充建议','COLLABORATION')

# 协同机制
y=800000
tb(s,M+200000,y,5000000,300000,'协同机制',fs=20,c=BLUE,b=True)

collabs=[
    ('周度','推广与运营同步异常+动作','推广助理输出周度异常扫描，负责人标注优先级，同步给对应运营'),
    ('月度','推广输出品类增长建议','基于月度数据复盘，输出品类增长建议表，指导运营调整商品策略'),
    ('大促前','四部门联审','推广+运营+商品+供应链：货品池、预算池、资源位、页面承接联审，避免流量到位但承接不足'),
]
for i,(freq,act,desc) in enumerate(collabs):
    cy=y+400000+i*700000
    rect(s,M,cy,SW-M*2,600000,fill=CARD,r=80000)
    colors_c=[BLUE,GREEN,AMBER]
    rect(s,M,cy,45000,600000,fill=colors_c[i])
    tb(s,M+150000,cy+50000,2000000,220000,freq,fs=15,c=colors_c[i],b=True)
    tb(s,M+150000,cy+270000,SW-M*2-300000,160000,act,fs=12,c=WHITE,b=True)
    tb(s,M+150000,cy+430000,SW-M*2-300000,140000,desc,fs=10,c=GREY)

# 补充建议
y2=3400000
tb(s,M+200000,y2,5000000,300000,'补充建议',fs=20,c=PURPLE,b=True)

extras=[
    ('建立推广知识库','将商品分层SOP、异常处理清单、AI指令库、复盘模板统一归档，新人可快速上手',PURPLE),
    ('跨部门数据打通','推动推广数据与商品库存、供应链、财务数据的联动，让推广决策有更完整的经营视角',BLUE),
    ('季度推广策略会','每季度组织一次推广策略复盘与规划会，推广+运营+品类共同参与，形成策略共识',GREEN),
]
for i,(t,desc,cl) in enumerate(extras):
    ey=y2+400000+i*600000
    rect(s,M,ey,SW-M*2,520000,fill=CARD,r=80000)
    rect(s,M,ey,40000,520000,fill=cl)
    tb(s,M+130000,ey+50000,SW-M*2-260000,200000,t,fs=13,c=cl,b=True)
    tb(s,M+130000,ey+270000,SW-M*2-260000,200000,desc,fs=10,c=GREY)

# ═══════════════════════════════════
# SLIDE 11: 致谢
# ═══════════════════════════════════
s=prs.slides.add_slide(BL); bg(s)
tb(s,0,2200000,SW,750000,'THANK YOU',fs=54,c=WHITE,b=True,al=PP_ALIGN.CENTER)
rect(s,(SW-3800000)//2,3100000,3800000,22000,fill=BLUE)
tb(s,0,3350000,SW,380000,'让增长更可控，让效率更稳定',fs=17,c=GREY,al=PP_ALIGN.CENTER)
tb(s,0,4700000,SW,280000,'国内营销中心  |  东方  |  2026.07.23',fs=12,c=GREY2,al=PP_ALIGN.CENTER)

# ═══ 保存 ═══
out='/Users/llano/Desktop/国内hermes存储/东方/2026推广年中会议汇报_专家版.pptx'
prs.save(out)
print(f'✅ {out}')
print(f'{len(prs.slides)} 页')
