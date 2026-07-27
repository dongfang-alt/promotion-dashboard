#!/usr/bin/env python3
"""
2026推广年中会议汇报 — 专家深度改写版
基于原PPT逐页重写：每句话都有判断，每个数字都有归因
"""
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ═══ 设计系统（保持深色商务风） ═══
BG     =RGBColor(0x0A,0x0F,0x1E); CARD=RGBColor(0x14,0x1E,0x33)
HEADER =RGBColor(0x0F,0x17,0x2A); ACCENT_BG=RGBColor(0x1E,0x29,0x3B)
BLUE   =RGBColor(0x3B,0x82,0xF6); BLUE_BG=RGBColor(0x1E,0x3A,0x5F)
GREEN  =RGBColor(0x22,0xC5,0x5E); GREEN_BG=RGBColor(0x14,0x52,0x2D)
AMBER  =RGBColor(0xF5,0x9E,0x0B); AMBER_BG=RGBColor(0x78,0x3E,0x0F)
RED    =RGBColor(0xEF,0x44,0x44); RED_BG  =RGBColor(0x7F,0x1D,0x1D)
PURPLE =RGBColor(0x8B,0x5C,0xF6); PURPLE_BG=RGBColor(0x3B,0x1F,0x6E)
WHITE  =RGBColor(0xF8,0xFA,0xFC); GREY=RGBColor(0x94,0xA3,0xB8)
GREY2  =RGBColor(0x64,0x74,0x8B); BORDER_C=RGBColor(0x33,0x40,0x55)

SW=12192000; SH=6858000; M=720000; G=180000

prs=Presentation(); prs.slide_width=SW; prs.slide_height=SH
BL=prs.slide_layouts[6]

# ═══ 工具函数 ═══
def bg(s,c=BG):
    r=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,SW,SH)
    r.fill.solid();r.fill.fore_color.rgb=c;r.line.fill.background()

def rc(s,l,t,w,h,fill=None,bc=None,bw=None,rd=None):
    shape=MSO_SHAPE.ROUNDED_RECTANGLE if rd else MSO_SHAPE.RECTANGLE
    r=s.shapes.add_shape(shape,l,t,w,h)
    if fill:r.fill.solid();r.fill.fore_color.rgb=fill
    else:r.fill.background()
    if bc:r.line.color.rgb=bc
    if bw:r.line.width=bw
    else:r.line.fill.background()
    return r

def tb(s,l,t,w,h,text,fs=14,c=WHITE,b=False,al=PP_ALIGN.LEFT,fn='微软雅黑',ls=1.2):
    tx=s.shapes.add_textbox(l,t,w,h);tf=tx.text_frame;tf.word_wrap=True
    p=tf.paragraphs[0];p.text=text
    p.font.size=Pt(fs);p.font.color.rgb=c;p.font.bold=b;p.font.name=fn
    p.alignment=al;p.space_after=Pt(1);p.line_spacing=Pt(fs*ls)
    return tx

def mtb(s,l,t,w,h,lines,fs=12,c=GREY,fn='微软雅黑',ls=1.25):
    tx=s.shapes.add_textbox(l,t,w,h);tf=tx.text_frame;tf.word_wrap=True
    for i,(text,opts) in enumerate(lines):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        sz=opts.get('size',fs);cl=opts.get('color',c);bo=opts.get('bold',False)
        p.text=text;p.font.size=Pt(sz);p.font.color.rgb=cl;p.font.bold=bo
        p.font.name=fn;p.space_after=Pt(2);p.line_spacing=Pt(sz*ls)
    return tx

def hdr(s,title,sub=""):
    rc(s,0,0,SW,550000,fill=HEADER);rc(s,0,550000,SW,20000,fill=BLUE)
    tb(s,M,110000,7000000,330000,title,fs=24,c=WHITE,b=True)
    if sub:tb(s,SW-M-3800000,140000,3800000,250000,sub,fs=11,c=GREY,al=PP_ALIGN.RIGHT)

def kpi_card(s,l,t,w,label,value,vc=None,sub="",bg_c=CARD,ls=9,vs=14):
    rc(s,l,t,w,500000,fill=bg_c,rd=60000)
    tb(s,l+60000,t+45000,w-120000,130000,label,fs=ls,c=GREY)
    tb(s,l+60000,t+175000,w-120000,200000,value,fs=vs,c=vc or WHITE,b=True)
    if sub:tb(s,l+60000,t+370000,w-120000,110000,sub,fs=8,c=GREY)

def div(s,l,t,w):rc(s,l,t,w,10000,fill=BORDER_C)

# ═══════════════════════════════════════
# SLIDE 1: 封面
# ═══════════════════════════════════════
s1=prs.slides.add_slide(BL);bg(s1)
rc(s1,0,0,SW,40000,fill=BLUE)
rc(s1,M,1200000,35000,2400000,fill=BLUE)
tb(s1,M+160000,1200000,9500000,750000,'2026 推广年中会议汇报',fs=44,c=WHITE,b=True)
tb(s1,M+160000,1950000,8000000,300000,'E-commerce Promotion Mid-Year Review',fs=16,c=GREY)
tb(s1,M+160000,2350000,7000000,320000,'H1 平台复盘  /  H2 三大规划  /  组织升级',fs=20,c=BLUE,b=True)
div(s1,M+160000,2770000,6000000)
y0=4850000
tb(s1,M+160000,y0,3000000,260000,'汇报人：东方',fs=13,c=GREY)
tb(s1,M+160000,y0+310000,3000000,260000,'部门：国内营销中心',fs=13,c=GREY)
tb(s1,M+160000,y0+620000,3000000,260000,'汇报时间：2026.07.23',fs=13,c=GREY)
tb(s1,SW-M-3800000,SH-M-300000,3800000,300000,'从投放执行到增长经营',fs=15,c=BLUE,b=True,al=PP_ALIGN.RIGHT)

# ═══════════════════════════════════════
# SLIDE 2: 目录
# ═══════════════════════════════════════
s2=prs.slides.add_slide(BL);bg(s2)
hdr(s2,'汇报结构','CONTENTS')

mods=[('01','上半年复盘','天猫旗舰店与京东自营的\n效率修复与增长矛盾',BLUE,BLUE_BG),
      ('02','下半年规划','商品分层精细化运营\n品类协同与渠道差异化',GREEN,GREEN_BG),
      ('03','组织升级','降低个人经验依赖\nAI工具化与团队能力建设',PURPLE,PURPLE_BG)]
cw_m=3400000;ch_m=2300000;cg_m=250000
tw_m=cw_m*3+cg_m*2;sx_m=(SW-tw_m)//2;y_m=900000
for i,(num,title,desc,cl,bc) in enumerate(mods):
    x=sx_m+i*(cw_m+cg_m)
    rc(s2,x,y_m,cw_m,ch_m,fill=bc,rd=120000);rc(s2,x,y_m,cw_m,50000,fill=cl)
    tb(s2,x+200000,y_m+150000,cw_m-400000,400000,num,fs=42,c=cl,b=True)
    tb(s2,x+200000,y_m+520000,cw_m-400000,350000,title,fs=22,c=WHITE,b=True)
    tb(s2,x+200000,y_m+950000,cw_m-400000,1000000,desc,fs=13,c=GREY)

ly=y_m+ch_m+350000
rc(s2,sx_m,ly,tw_m,750000,fill=ACCENT_BG,rd=100000)
tb(s2,sx_m+250000,ly+120000,2200000,280000,'汇报主线',fs=17,c=BLUE,b=True)
tb(s2,sx_m+250000,ly+420000,tw_m-500000,280000,'H1完成推广效率全面修复，H2从「止血」转向「造血」——精细化运营、品类协同、渠道差异化，推动推广从执行动作升级为增长经营能力。',fs=12,c=GREY)

# ═══════════════════════════════════════
# SLIDE 3: 天猫旗舰店 — H1增长账本
# ═══════════════════════════════════════
s3=prs.slides.add_slide(BL);bg(s3)
hdr(s3,'上半年复盘 01：天猫旗舰店','消耗+30.3%，ROI同比+39.3%——钱花得更值了')

tb(s3,M,640000,10600000,240000,'同比=2026H1 vs 2025H1 | 环比=2026H1 vs 2025H2 | ROI用数值，其他用百分比',fs=10,c=GREY2)

# KPI行
y3=930000;lw3=1700000;kw3=1060000
rc(s3,M,y3,lw3,1020000,fill=BLUE_BG,rd=80000)
tb(s3,M+100000,y3+130000,lw3-200000,240000,'天猫旗舰店',fs=15,c=WHITE,b=True)
tb(s3,M+100000,y3+430000,lw3-200000,200000,'核心店铺',fs=10,c=BLUE,b=True)
tb(s3,M+100000,y3+660000,lw3-200000,280000,'H1主线：效率修复+销售增长',fs=9,c=GREY)

tm_kpi=[("消耗同比","+30.3%",AMBER),("消耗环比","-0.8%",GREY),("ROI同比","+39.3%",GREEN),
        ("ROI环比","+10.4%",GREEN),("费比同比","+2.6pt",AMBER),("费比环比","-0.9pt",GREEN),
        ("销售同比","+10.0%",GREEN),("销售环比","+4.4%",GREEN)]
for i,(lb,vl,vc) in enumerate(tm_kpi):
    x=M+lw3+50000+i*(kw3+35000);kpi_card(s3,x,y3,kw3,lb,vl,vc)

# 数码旗舰
y3b=2050000
rc(s3,M,y3b,lw3,750000,fill=ACCENT_BG,rd=80000)
tb(s3,M+100000,y3b+130000,lw3-200000,220000,'天猫数码旗舰店',fs=14,c=WHITE,b=True)
tb(s3,M+100000,y3b+400000,lw3-200000,250000,'新起盘 | 轻预算测试 | 先建规则再扩量',fs=9,c=GREY)
dig_lbl=["消耗同比","消耗环比","ROI同比","ROI环比","费比同比","费比环比","CPC同比","CVR同比"]
for i,lb in enumerate(dig_lbl):
    x=M+lw3+50000+i*(kw3+35000)
    kpi_card(s3,x,y3b,kw3,lb,'新起盘',vc=GREY,bg_c=ACCENT_BG,ls=8,vs=11)

# 核心结论三卡
y3c=3000000;ccw=3200000;ccg_c=250000;ccs=M+400000
ccs_d=[('ROI修复','ROI环比+10.4%','H1推广优化效果明确——消耗放量的同时，每一块钱的效率在提升',GREEN,GREEN_BG),
       ('费比改善','环比-0.9pt','同销售口径下费用效率在改善，控费与提效同步成立',GREEN,GREEN_BG),
       ('Q2联合投放','销售同比+9.6%/环比+24.4%','Q2签订平台联合高费比合同，在ROI和CPC可控下滑的前提下，换来了销售确定性增长。这是H2可复制的模型',BLUE,BLUE_BG)]
for i,(t,v,d,cl,bc) in enumerate(ccs_d):
    x=ccs+i*(ccw+ccg_c)
    rc(s3,x,y3c,ccw,850000,fill=bc,rd=100000)
    tb(s3,x+180000,y3c+100000,ccw-360000,180000,t,fs=11,c=GREY)
    tb(s3,x+180000,y3c+300000,ccw-360000,220000,v,fs=17,c=cl,b=True)
    tb(s3,x+180000,y3c+540000,ccw-360000,260000,d,fs=9,c=GREY)

# 诊断+数码
y3d=4100000;dw=5300000
rc(s3,M,y3d,dw,1200000,fill=CARD,rd=100000)
rc(s3,M+150000,y3d+80000,40000,250000,fill=BLUE)
tb(s3,M+250000,y3d+70000,dw-400000,260000,'天猫旗舰店核心判断',fs=16,c=BLUE,b=True)
mtb(s3,M+250000,y3d+380000,dw-400000,750000,[
    ('H1结论：天猫旗舰店是标杆——在合理费比框架内，规模与效率可以兼得。',{'size':12,'color':WHITE,'bold':True}),
    ('消耗+30.3%不是浪费，是策略性投入。Q2联合平台高费比合同已经验证：只要绑定活动目标、商品池和承接结果，加大投入能换来确定性的销售增长。',{}),
    ('H2方向：保爆款、控节奏、抢核心词，同时将Q2联合投放模型复制到更多活动节点。站外UD引流必须回看站内承接效率。',{}),
],fs=11,c=GREY)

rc(s3,M+dw+200000,y3d,dw+180000,1200000,fill=CARD,rd=100000)
rc(s3,M+dw+350000,y3d+80000,40000,250000,fill=PURPLE)
tb(s3,M+dw+450000,y3d+70000,dw-200000,260000,'天猫数码旗舰店',fs=16,c=PURPLE,b=True)
mtb(s3,M+dw+450000,y3d+380000,dw-200000,750000,[
    ('新起盘账户，缺少完整同比/环比评价基础——不急于下结论。',{'size':12,'color':PURPLE,'bold':True}),
    ('短期重点：模型冷启、人群沉淀、素材跑数、商品承接。H2以轻预算测试为主。',{}),
    ('长期定位：天猫平台的多店矩阵补充，先建规则再谈扩量。',{}),
],fs=11,c=GREY)

y3e=5600000
rc(s3,M,y3e,SW-M*2,450000,fill=ACCENT_BG,rd=80000)
tb(s3,M+200000,y3e+50000,SW-M*2-400000,350000,'天猫底线：旗舰店是主线——Q2已验证联合投放模型，H2复制+优化；数码旗舰店作为多店矩阵的战略补充，轻预算探路。',fs=12,c=WHITE,b=True)

# ═══════════════════════════════════════
# SLIDE 4: 京东自营 — 效率修复后的新课题
# ═══════════════════════════════════════
s4=prs.slides.add_slide(BL);bg(s4)
hdr(s4,'上半年复盘 02：京东自营','ROI环比+36.4%——H1效率修复最亮眼的数字')

tb(s4,M,640000,10600000,200000,'但销售同比-15.8%：控费的红利已经吃完，继续收缩就会伤及基本盘。',fs=12,c=AMBER,b=True)

# KPI行
y4=950000
rc(s4,M,y4,lw3,1020000,fill=AMBER_BG,rd=80000)
tb(s4,M+100000,y4+130000,lw3-200000,240000,'京东自营',fs=15,c=WHITE,b=True)
tb(s4,M+100000,y4+430000,lw3-200000,200000,'核心店铺',fs=10,c=AMBER,b=True)
tb(s4,M+100000,y4+660000,lw3-200000,280000,'H1主线：效率大幅修复，规模承压',fs=9,c=GREY)

jd_kpi=[("消耗同比","-13.2%",GREY),("消耗环比","-15.8%",GREY),("ROI同比","+25.3%",GREEN),
        ("ROI环比","+36.4%",GREEN),("费比同比","+0.3pt",AMBER),("费比环比","-1.0pt",GREEN),
        ("销售同比","-15.8%",RED),("销售环比","-8.1%",RED)]
for i,(lb,vl,vc) in enumerate(jd_kpi):
    x=M+lw3+50000+i*(kw3+35000);kpi_card(s4,x,y4,kw3,lb,vl,vc)

# POP + MrGreen
y4b=2070000;sub_kw=950000
# POP
rc(s4,M,y4b,lw3,750000,fill=ACCENT_BG,rd=80000)
tb(s4,M+100000,y4b+130000,lw3-200000,220000,'京东POP',fs=14,c=WHITE,b=True)
tb(s4,M+100000,y4b+400000,lw3-200000,250000,'ROI+35.1% | 结构补充',fs=10,c=GREEN,b=True)
for i,(lb,vl) in enumerate(zip(dig_lbl,["+3.2%","-10.8%","+26.3%","+35.1%","+0.1pt","-0.6pt","-17.9%","+12.9%"])):
    x=M+lw3+50000+i*(sub_kw+28000)
    vc_c=GREEN if ('ROI' in lb or 'CVR' in lb) else GREY
    kpi_card(s4,x,y4b,sub_kw,lb,vl,vc=vc_c,bg_c=ACCENT_BG,ls=8,vs=11)

# MrGreen
y4c=2920000
rc(s4,M,y4c,lw3,750000,fill=ACCENT_BG,rd=80000)
tb(s4,M+100000,y4c+130000,lw3-200000,220000,'Mr.Green 自营',fs=14,c=WHITE,b=True)
tb(s4,M+100000,y4c+400000,lw3-200000,250000,'ROI+8.5% | 利润型补充',fs=10,c=GREEN,b=True)
for i,(lb,vl) in enumerate(zip(dig_lbl,["-19.3%","-7.2%","+8.5%","+6.5%","-0.9pt","-0.7pt","-12.2%","+2.1%"])):
    x=M+lw3+50000+i*(sub_kw+28000)
    vc_c=GREEN if ('ROI' in lb or 'CVR' in lb) else GREY
    kpi_card(s4,x,y4c,sub_kw,lb,vl,vc=vc_c,bg_c=ACCENT_BG,ls=8,vs=11)

# 三卡结论
y4d=3850000
jd_ccs=[('ROI修复','ROI环比+36.4%','H1最亮眼的效率数据——自营修复最为显著，控费+提效同步成立',GREEN,GREEN_BG),
        ('费比改善','环比-1.0pt','同销售费比明显下降，费用效率在改善，控费策略有效',GREEN,GREEN_BG),
        ('规模承压','销售同比-15.8%/环比-8.1%','效率修复的代价——过度控费压制了增长。H2的核心命题是在守住ROI底线的同时恢复扩量',AMBER,AMBER_BG)]
for i,(t,v,d,cl,bc) in enumerate(jd_ccs):
    x=ccs+i*(ccw+ccg_c)
    rc(s4,x,y4d,ccw,850000,fill=bc,rd=100000)
    tb(s4,x+180000,y4d+100000,ccw-360000,180000,t,fs=11,c=GREY)
    tb(s4,x+180000,y4d+300000,ccw-360000,220000,v,fs=17,c=cl,b=True)
    tb(s4,x+180000,y4d+540000,ccw-360000,260000,d,fs=9,c=GREY)

# 诊断
y4e=4950000
rc(s4,M,y4e,dw,1200000,fill=CARD,rd=100000)
rc(s4,M+150000,y4e+80000,40000,250000,fill=AMBER)
tb(s4,M+250000,y4e+70000,dw-400000,260000,'京东自营核心判断',fs=16,c=AMBER,b=True)
mtb(s4,M+250000,y4e+380000,dw-400000,750000,[
    ('H1的成绩：ROI修复——自营环比+36.4%，POP环比+35.1%，效果非常明确。',{'size':12,'color':WHITE,'bold':True}),
    ('H1的代价：销售同比-15.8%——消耗砍了15.8%，规模必然收缩。控费策略有效，但已接近临界点。',{}),
    ('H2关键选择：不是要不要控费的问题，而是在什么ROI底线之上恢复扩量。关键词质量分、人群分层、爆品预算保护三管齐下。',{}),
],fs=11,c=GREY)

rc(s4,M+dw+200000,y4e,dw+180000,1200000,fill=CARD,rd=100000)
rc(s4,M+dw+350000,y4e+80000,40000,250000,fill=PURPLE)
tb(s4,M+dw+450000,y4e+70000,dw-200000,260000,'POP与Mr.Green',fs=16,c=PURPLE,b=True)
mtb(s4,M+dw+450000,y4e+380000,dw-200000,750000,[
    ('POP：ROI+35.1%，效率修复与自营同步。策略上按转化稳定性筛选投入，不追求全面铺量。',{'size':12,'color':PURPLE,'bold':True}),
    ('Mr.Green：ROI+8.5%，费比双降，可承担利润型补充角色。小众品类，稳定贡献即可。',{}),
    ('两者不做主叙事，作为京东平台结构的必要组成。',{}),
],fs=11,c=GREY)

y4f=6400000
rc(s4,M,y4f,SW-M*2,400000,fill=ACCENT_BG,rd=80000)
tb(s4,M+200000,y4f+45000,SW-M*2-400000,300000,'京东底线：H1效率修复已完成，H2必须恢复增长弹性——在守住ROI底线的同时，让核心品类和爆品重新扩量。',fs=12,c=WHITE,b=True)

# ═══════════════════════════════════════
# SLIDE 5: H2-01 推广精细化提效
# ═══════════════════════════════════════
s5=prs.slides.add_slide(BL);bg(s5)
hdr(s5,'下半年规划 01：推广精细化提效','从一个人盯核心品，到系统盯全盘')

tb(s5,M,640000,10600000,280000,'背景：H1一人统筹多店铺多模式，精力只能盯核心品——这是H1三个压力点的根源。H2通过推广助理+AI，把监控范围从「核心品」扩展到「全品类×全计划」。',fs=11,c=GREY)

# 四层级
y5=1000000;qw=2620000;qh=1800000
quads=[('核心爆款',BLUE,BLUE_BG,'保排名·保转化·保预算','看ROI环比/核心词排名/活动承接','消耗占60%+，优先级最高'),
       ('潜力单品',GREEN,GREEN_BG,'测人群·测素材·测场景','看CPC/CVR/加购率/ROI趋势','每周评估是否升级为核心爆款'),
       ('新品',AMBER,AMBER_BG,'小预算冷启·先跑模型','看点击率/加购率/素材胜出率','冷启周期≤14天，到期评估'),
       ('中长尾品',PURPLE,PURPLE_BG,'低预算验证·快速止损','看ROI阈值/转化率/费用泄露','连续7天低于盈亏线→自动暂停')]
for i,(t,cl,bc,act,watch,rule) in enumerate(quads):
    x=M+i*(qw+170000)
    rc(s5,x,y5,qw,qh,fill=bc,rd=100000);rc(s5,x,y5,qw,40000,fill=cl)
    tb(s5,x+150000,y5+120000,qw-300000,260000,t,fs=16,c=cl,b=True)
    tb(s5,x+150000,y5+420000,qw-300000,200000,act,fs=12,c=WHITE,b=True)
    tb(s5,x+150000,y5+660000,qw-300000,400000,watch,fs=10,c=GREY)
    tb(s5,x+150000,y5+1100000,qw-300000,300000,rule,fs=9,c=GREY2)

# 三列升级
y5b=3000000;c3w=3400000
ups=[('监控范围升级',BLUE,BLUE_BG,['过去：核心品类+爆款+主要账户','升级：平台→店铺→品类→商品→计划→关键词→人群→素材','目标：从事后发现大问题→过程识别小异常']),
     ('人力分工升级',GREEN,GREEN_BG,['负责人：策略/预算/平台打法/关键决策','助理：数据整理/异常初筛/计划跟进','目标：负责人从拉表中释放，投入策略']),
     ('AI工具提效',PURPLE,PURPLE_BG,['AI：数据处理+异常识别+趋势归因+策略草案','人工：业务判断+动作取舍+执行+复盘','目标：提速诊断，降低个人经验依赖'])]
c3s=(SW-c3w*3-200000*2)//2
for i,(t,cl,bc,items) in enumerate(ups):
    x=c3s+i*(c3w+200000)
    rc(s5,x,y5b,c3w,1750000,fill=bc,rd=120000);rc(s5,x,y5b,c3w,40000,fill=cl)
    tb(s5,x+180000,y5b+120000,c3w-360000,260000,t,fs=15,c=cl,b=True)
    ils=[(f'▸ {it}',{'size':10,'color':GREY}) for it in items]
    mtb(s5,x+180000,y5b+450000,c3w-360000,1200000,ils,fs=10,c=GREY,ls=1.3)

# 落地
y5c=5000000
rc(s5,M,y5c,SW-M*2,420000,fill=ACCENT_BG,rd=80000)
tb(s5,M+200000,y5c+40000,SW-M*2-400000,340000,'落地机制：日看异常（AI自动扫描→推送）→周看商品层级迁移（助理整理→负责人判断）→月看平台结构（负责人输出策略报告）。核心原则：80%精力给20%核心爆款，系统覆盖长尾。',fs=11,c=WHITE,b=True)

# ═══════════════════════════════════════
# SLIDE 6: H2-02 赋能品类增长
# ═══════════════════════════════════════
s6=prs.slides.add_slide(BL);bg(s6)
hdr(s6,'下半年规划 02：赋能品类增长','推广的终局不是买流量，是帮品类找到增长路径')

tb(s6,M,640000,10600000,280000,'目标：通过投放数据反向输出经营建议——哪些品类值得加预算、哪些商品需要调价格、哪些页面承接有问题。让推广从成本中心变为增长参谋。',fs=12,c=GREY)

y6=1050000;bcw=3400000;bch=3800000
bcs=(SW-bcw*3-200000*2)//2
bigs=[('数据反馈：识别机会',GREEN,
      ['按品类×商品监控ROI/CPC/CVR/点击率/收藏加购，建立趋势基线',
       '三类机会自动标记：低CPC高CVR→建议放量；高点击低转化→承接问题，同步运营；高消耗低ROI→建议止损',
       '输出给运营的具体动作：加预算的品类清单、需优化页面/价格/权益的商品清单']),
     ('组合打法：提升效率',BLUE,
      ['搜索=承接明确需求（保转化）；推荐/场景=拓新人群（扩漏斗）；站外种草=补认知（蓄水）',
       '核心品：搜索+推荐联动，保排名保转化；潜力品：小预算多场景测试；新品：先内容种草再站内承接',
       '每类打法绑定复盘指标，不看动作看结果——花了多少钱、带来了多少销售、ROI是否在改善']),
     ('经营协同：放大结果',PURPLE,
      ['推广结论同步商品/运营/供应链：爆款备货提醒、价格竞争力分析、详情页转化率诊断',
       '大促前：推广+运营+商品+供应链联审货品池/预算池/资源位/页面承接——避免流量到了但库存/权益/承接跟不上',
       '沉淀品类推广打法库，形成「什么品类在什么阶段用什么打法」的可复用模型'])]
for i,(t,cl,items) in enumerate(bigs):
    x=bcs+i*(bcw+200000)
    rc(s6,x,y6,bcw,bch,fill=CARD,rd=120000);rc(s6,x,y6,bcw,45000,fill=cl)
    tb(s6,x+220000,y6+130000,bcw-440000,280000,t,fs=17,c=cl,b=True)
    ils=[(f'{j+1}. {it}',{'size':10,'color':GREY}) for j,it in enumerate(items)]
    mtb(s6,x+220000,y6+500000,bcw-440000,3100000,ils,fs=10,c=GREY,ls=1.3)

y6b=5150000
rc(s6,M,y6b,SW-M*2,420000,fill=ACCENT_BG,rd=80000)
tb(s6,M+200000,y6b+40000,SW-M*2-400000,340000,'核心输出：不是交一份推广复盘表，而是输出一份品类增长建议表——让推广数据真正成为经营决策的依据。',fs=13,c=WHITE,b=True)

# ═══════════════════════════════════════
# SLIDE 7: H2-03 渠道策略
# ═══════════════════════════════════════
s7=prs.slides.add_slide(BL);bg(s7)
hdr(s7,'下半年规划 03：渠道策略方向','天猫 / 京东 / 拼多多')

tb(s7,M,640000,10600000,260000,'核心逻辑：不再用同一套ROI标准一刀切——天猫用可控费比换增长，京东在效率底线上恢复规模，拼多多最小成本跑通方法。',fs=12,c=GREY)

y7=1050000;chw=3400000;chh=4400000
chs=(SW-chw*3-200000*2)//2
chs_d=[('天猫','可控费比换增长',BLUE,BLUE_BG,
       ['定位：品牌主阵地，市场份额承接渠道',
        '策略：复制Q2联合投放模型——平台联合资源+活动节点+核心爆款，以可控高费比换取确定性销售增长',
        '动作：核心词卡位、推荐扩量、UD站外引流、内容种草回流',
        '边界：高费比必须绑定活动目标+商品池+承接结果。Q2已验证了模型，H2是复制+优化，不是无目标放量',
        '关键指标：销售增速 vs ROI下滑幅度——不能只看ROI，也不能只看销售']),
      ('京东','在效率底线上恢复规模',AMBER,AMBER_BG,
       ['定位：效率修复已完成，H2进入效率+规模平衡阶段',
        '策略：以自营为主，守住高ROI计划不放松，同时核心品类和爆品恢复扩量——不是不控费，是换一个更高的费比目标来控',
        '动作：关键词质量分优化、人群分层扩量、低效词清理、爆品预算保护',
        '边界：优先看ROI环比/费比环比/CPC/CVR变化。如果扩量后ROI环比下滑超过X%，回调预算节奏',
        '关键矛盾：H1销售-16%不可持续——不补量，品类市场份额会继续被蚕食']),
      ('拼多多','最小成本跑通方法',PURPLE,PURPLE_BG,
       ['定位：9月接手，新增提效场——前两个月不追ROI',
        '策略：先完成账户体检+商品分层+计划结构+投放规则搭建，再进入规模化提效',
        '动作：搜索/场景基础计划搭建、可投商品池筛选、素材AB测试、ROI阈值+止损机制设置',
        '边界：9-10月以学习和验证为主，双11小规模参与验证，12月沉淀规则后做2027年规划',
        '节奏：9月诊断→10月测试→11月小规模验证→12月定规则'])]
for i,(name,sub,cl,bc,items) in enumerate(chs_d):
    x=chs+i*(chw+200000)
    rc(s7,x,y7,chw,chh,fill=bc,rd=120000);rc(s7,x,y7,chw,50000,fill=cl)
    tb(s7,x+200000,y7+120000,chw-400000,280000,name,fs=22,c=cl,b=True)
    tb(s7,x+200000,y7+400000,chw-400000,240000,sub,fs=13,c=WHITE,b=True)
    ils=[(f'▸ {it}',{'size':10,'color':GREY}) for it in items]
    mtb(s7,x+200000,y7+750000,chw-400000,3500000,ils,fs=10,c=GREY,ls=1.25)

y7b=5700000
rc(s7,M,y7b,SW-M*2,420000,fill=ACCENT_BG,rd=80000)
tb(s7,M+200000,y7b+40000,SW-M*2-400000,340000,'渠道判断：天猫用可控费比换确定性增长（Q2模型已验证）→京东在效率底线上恢复规模弹性（控费已到临界点）→拼多多先建规则再规模化（9月接手，不急着追ROI）。',fs=12,c=WHITE,b=True)

# ═══════════════════════════════════════
# SLIDE 8: 组织01 — 成员管理
# ═══════════════════════════════════════
s8=prs.slides.add_slide(BL);bg(s8)
hdr(s8,'组织升级 01：组织成员管理','把个人经验拆成团队方法，降低对关键人员的依赖')

tb(s8,M,640000,10600000,260000,'核心命题：H1一个人盯全局的瓶颈已经显现（长尾品失控、数据工具缺失、全链路断层）。H2需要从「个人能力驱动」升级为「团队+系统驱动」。',fs=12,c=GREY)

y8=1050000;ocw=3400000;och=3800000
ocs=(SW-ocw*3-200000*2)//2
orgs=[('能力建设：从操作到判断',BLUE,
       ['现状：成员能做基础执行，但缺少平台机制、商品层级、预算节奏的独立判断能力',
        '方向：从「会建计划、会调价」→「会看趋势、会归因、会判断扩量/止损」',
        '方法：ROI/费比/CPC/CVR联动分析训练，商品池分层实操，活动与推广节奏匹配演练',
        '案例：铁蛋半年培养从部门助理→推广助理→成功转岗，验证方法论可复制']),
      ('分工机制：负责人+助理协同',GREEN,
       ['负责人：平台策略、预算分配、关键账户诊断、跨部门协同——把时间花在判断上',
        '助理：日报数据、异常初筛、计划跟进、素材/关键词维护——把重复劳动标准化',
        '协同SOP：助理每天早上输出昨日异常扫描→负责人10分钟内标注优先级→当天给动作',
        '铁蛋加入后将承担天猫旗舰店的日报监控和异常初筛，释放负责人做策略判断的时间']),
      ('标准沉淀：方法论产品化',PURPLE,
       ['商品分层SOP：核心爆款/潜力/新品/中长尾→对应投放规则+预算比例+止损线',
        '异常处理清单：ROI下滑/费比上升/CPC异常/CVR波动/消耗突增→对应排查路径+处置动作',
        '复盘机制：周看计划级异常→月看平台结构→活动看商品+承接→季度输出方法论迭代',
        '所有SOP和清单统一归档到推广知识库，新人可在一周内上手基础监控'])]
for i,(t,cl,items) in enumerate(orgs):
    x=ocs+i*(ocw+200000)
    rc(s8,x,y8,ocw,och,fill=CARD,rd=120000);rc(s8,x,y8,ocw,45000,fill=cl)
    tb(s8,x+200000,y8+130000,ocw-400000,280000,t,fs=16,c=cl,b=True)
    ils=[(f'▸ {it}',{'size':10,'color':GREY}) for it in items]
    mtb(s8,x+200000,y8+500000,ocw-400000,3100000,ils,fs=10,c=GREY,ls=1.2)

y8b=5150000
rc(s8,M,y8b,SW-M*2,420000,fill=ACCENT_BG,rd=80000)
tb(s8,M+200000,y8b+40000,SW-M*2-400000,340000,'核心逻辑：铁蛋的半年培养→转岗证明了方法论可复制。H2把这个模式系统化——用SOP降低上手门槛，用分工释放负责人精力，用复盘持续迭代方法。',fs=12,c=WHITE,b=True)

# ═══════════════════════════════════════
# SLIDE 9: 组织02 — AI赋能
# ═══════════════════════════════════════
s9=prs.slides.add_slide(BL);bg(s9)
hdr(s9,'组织升级 02：AI 赋能团队','AI不是替代判断，是让判断更快、更准、更标准化')

tb(s9,M,640000,10600000,280000,'核心逻辑：AI做数据处理+异常识别+策略草案——把推广负责人从拉表和排查中释放出来。人工做业务判断+动作取舍+执行落地+复盘纠偏——这是AI替代不了的。',fs=12,c=GREY)

y9=1100000
ais=[('AI多维数据处理',PURPLE,
     ['每天自动扫描全平台×全店铺×全品类×全计划的数据，识别消耗/ROI/CPC/CVR的异常波动',
      '从「人工拉表找问题」（周级、易遗漏）→「AI先筛异常，人工判断原因」（日级、全覆盖）',
      '具体场景：每天早上9:00，AI自动输出昨日异常扫描报告，推送到推广群，负责人10分钟标注优先级']),
     ('AI快速定位核心问题',BLUE,
      ['由总到细的下钻逻辑：平台概览→店铺对比→品类拆解→商品诊断→计划/关键词归因',
      '核心指标模块化：消耗、ROI、CPC、CVR、点击率、加购率、转化率——异常自动标记+趋势预警',
      'AI输出的不是原始数据，是「可能原因」：流量成本变化/转化承接问题/预算节奏失衡/商品竞争力下降/活动资源不足']),
     ('AI专业策略输出',GREEN,
      ['沉淀推广指令库：每种异常场景→对应的分析指令+优化建议模板',
      '运营提出组合问题（如：「XX品类最近CPC上升但CVR没变，怎么办」）→AI生成策略草案→负责人校验→指导实操',
      '价值：降低推广指导成本，让团队在没有资深推广在场时也能快速获得标准化判断'])]
for i,(t,cl,items) in enumerate(ais):
    x=ocs+i*(ocw+200000)
    rc(s9,x,y9,ocw,3400000,fill=CARD,rd=120000);rc(s9,x,y9,ocw,45000,fill=cl)
    tb(s9,x+200000,y9+130000,ocw-400000,280000,t,fs=15,c=cl,b=True)
    ils=[(f'▸ {it}',{'size':10,'color':GREY}) for it in items]
    mtb(s9,x+200000,y9+500000,ocw-400000,2700000,ils,fs=10,c=GREY,ls=1.25)

y9b=4800000
rc(s9,M,y9b,SW-M*2,420000,fill=ACCENT_BG,rd=80000)
tb(s9,M+200000,y9b+40000,SW-M*2-400000,340000,'AI闭环：每日数据清洗→异常自动识别→策略草案生成→人工判断+动作→执行→复盘→反哺指令库。关键原则：AI出初稿，人工做终判。',fs=12,c=WHITE,b=True)

# ═══════════════════════════════════════
# SLIDE 10: 行动节奏+协同
# ═══════════════════════════════════════
s10=prs.slides.add_slide(BL);bg(s10)
hdr(s10,'H2 行动节奏与协同机制','MILESTONE & COLLABORATION')

tb(s10,M,640000,10600000,240000,'H2主线：先建规则 → 接渠道 → 打大促 → 沉淀方法。每个阶段有可验收的产出。',fs=12,c=GREY)

y10=1000000;tlw=2550000;tlh=2000000;tlg=170000
tls=(SW-tlw*4-tlg*3)//2
tls_d=[('7-8月',BLUE,['产出① 商品分层SOP+预算分配规则','产出② 推广助理+AI日/周监控流程上线','产出③ 天猫/京东H2预算阈值+重点商品池锁定','产出④ 周度复盘模板+异常处理清单v1.0']),
       ('9月',PURPLE,['产出① 拼多多账户诊断报告','产出② 搜索/场景基础投放规则搭建','产出③ 双11预算框架+货品池初筛','产出④ 推广助理独立承担天猫日报监控']),
       ('10-11月',AMBER,['产出① 双11执行：预算节奏/爆品保护/止损','产出② 大促后复盘报告（平台×商品×素材）','产出③ 拼多多双11小规模验证报告','产出④ 品类增长建议表首次输出给运营']),
       ('12月',GREEN,['产出① 年度推广手册v1.0','产出② AI指令库+品类打法库','产出③ 2027年预算与渠道规划建议','产出④ 团队能力评估+下一年培养计划'])]
for i,(mo,cl,items) in enumerate(tls_d):
    x=tls+i*(tlw+tlg)
    rc(s10,x,y10,tlw,tlh,fill=CARD,rd=120000)
    rc(s10,x,y10,tlw,50000,fill=cl)
    tb(s10,x+180000,y10+140000,tlw-360000,280000,mo,fs=20,c=cl,b=True)
    ils=[(f'{it}',{'size':9,'color':GREY}) for it in items]
    mtb(s10,x+180000,y10+500000,tlw-360000,1400000,ils,fs=9,c=GREY,ls=1.25)

for i in range(3):
    ax=tls+(i+1)*tlw+i*tlg-150000
    tb(s10,ax,1850000,150000,300000,'→',fs=18,c=GREY2,b=True,al=PP_ALIGN.CENTER)

# 协同
y10b=3300000
rc(s10,M,y10b,SW-M*2,1500000,fill=CARD,rd=120000)
tb(s10,M+300000,y10b+100000,4000000,280000,'协同机制',fs=19,c=BLUE,b=True)

cols=[('周度','推广与运营同步','助理每周一输出上周异常扫描+本周关注清单，负责人标注优先级后同步对应运营',BLUE),
      ('月度','推广输出品类建议','基于月度数据复盘，输出品类增长建议表（含预算调整建议+商品优化建议），指导运营下月策略',GREEN),
      ('大促前','四部门联审','推广+运营+商品+供应链：货品池+预算池+资源位+页面承接联合审查，确保流量到达时全链路已就位',AMBER)]
for i,(fr,act,desc,cl) in enumerate(cols):
    iy=y10b+500000+i*300000
    rc(s10,M+350000,iy,35000,200000,fill=cl)
    tb(s10,M+500000,iy,1800000,200000,fr,fs=14,c=cl,b=True)
    tb(s10,M+2300000,iy,9000000,200000,f'{act}：{desc}',fs=11,c=GREY)

y10c=5100000
rc(s10,M,y10c,SW-M*2,420000,fill=ACCENT_BG,rd=80000)
tb(s10,M+200000,y10c+40000,SW-M*2-400000,340000,'年度目标：到2026年底，推广团队具备独立的数据诊断、策略输出和复盘沉淀能力——不再依赖单一个人的行业经验。',fs=12,c=WHITE,b=True)

# ═══════════════════════════════════════
# SLIDE 11: 致谢
# ═══════════════════════════════════════
s11=prs.slides.add_slide(BL);bg(s11)
tb(s11,0,2200000,SW,750000,'THANK YOU',fs=54,c=WHITE,b=True,al=PP_ALIGN.CENTER)
rc(s11,(SW-3800000)//2,3100000,3800000,22000,fill=BLUE)
tb(s11,0,3350000,SW,380000,'让增长更可控，让效率更稳定',fs=17,c=GREY,al=PP_ALIGN.CENTER)
tb(s11,0,4700000,SW,280000,'国内营销中心  |  东方  |  2026.07.23',fs=12,c=GREY2,al=PP_ALIGN.CENTER)

# ═══ 保存 ═══
out='/Users/llano/Desktop/国内hermes存储/东方/2026推广年中会议汇报_深度改写版.pptx'
prs.save(out)
print(f'✅ {out} | {len(prs.slides)}页')
