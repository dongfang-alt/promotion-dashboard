#!/usr/bin/env python3
"""
重建优化版 PPT — 2026推广年中会议汇报
设计系统：深色商务风，卡片式布局，KPI数据可视化
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import copy

# ═══════════════════════════════════════
# 设计系统
# ═══════════════════════════════════════

# 配色
BG_DARK      = RGBColor(0x0A, 0x0F, 0x1E)  # 最深背景
BG_CARD      = RGBColor(0x14, 0x1E, 0x33)  # 卡片背景
BG_HEADER    = RGBColor(0x0F, 0x17, 0x2A)  # 顶部栏
BG_ACCENT    = RGBColor(0x1E, 0x29, 0x3B)  # 次级卡片

BLUE         = RGBColor(0x3B, 0x82, 0xF6)  # 主蓝色
BLUE_LIGHT   = RGBColor(0x60, 0xA5, 0xFA)  # 浅蓝
BLUE_BG      = RGBColor(0x1E, 0x3A, 0x5F)  # 蓝色卡片背景
BLUE_BORDER  = RGBColor(0x25, 0x63, 0xEB)  # 蓝色边框

GREEN        = RGBColor(0x22, 0xC5, 0x5E)  # 正向绿
GREEN_BG     = RGBColor(0x14, 0x52, 0x2D)  # 绿色卡片
RED          = RGBColor(0xEF, 0x44, 0x44)  # 负向红
RED_BG       = RGBColor(0x7F, 0x1D, 0x1D)  # 红色卡片背景
AMBER        = RGBColor(0xF5, 0x9E, 0x0B)  # 琥珀色
AMBER_BG     = RGBColor(0x78, 0x3E, 0x0F)  # 琥珀底色
PURPLE       = RGBColor(0x8B, 0x5C, 0xF6)  # 紫色
PURPLE_BG    = RGBColor(0x3B, 0x1F, 0x6E)  # 紫色卡片
TEAL         = RGBColor(0x14, 0xB8, 0xA6)  # 青绿色

WHITE        = RGBColor(0xF8, 0xFA, 0xFC)  # 主文字
GREY         = RGBColor(0x94, 0xA3, 0xB8)  # 次要文字
GREY_DARK    = RGBColor(0x64, 0x74, 0x8B)  # 更暗文字
BORDER       = RGBColor(0x33, 0x40, 0x55)  # 分割线

# 尺寸
SLIDE_W = 12192000  # 13.33英寸
SLIDE_H = 6858000   # 7.5英寸
MARGIN  = 720000    # 0.79英寸边距
GAP     = 180000    # 卡片间距

# ═══════════════════════════════════════
# 工具函数
# ═══════════════════════════════════════

def add_bg(slide, color=BG_DARK):
    """添加全幅背景"""
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    return bg

def add_rect(slide, left, top, width, height, fill_color=None, border_color=None, border_width=None, radius=None):
    """添加矩形"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
                                    left, top, width, height)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    
    if border_color:
        shape.line.color.rgb = border_color
        if border_width:
            shape.line.width = border_width
    else:
        shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=14, color=WHITE, bold=False, 
                  alignment=PP_ALIGN.LEFT, font_name='微软雅黑', anchor=MSO_ANCHOR.TOP, line_spacing=1.2):
    """添加文本框"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(2)
    # line spacing
    p.line_spacing = Pt(font_size * line_spacing)
    
    # 设置东亚字体
    for run in p.runs:
        rPr = run._r.get_or_add_rPr()
        rPr.set('altLang', 'zh-CN')
    
    tf.auto_size = None
    return txBox

def add_multiline_text(slide, left, top, width, height, lines, font_size=13, color=WHITE, 
                        bold_first=False, font_name='微软雅黑', line_spacing=1.35):
    """添加多行文本框，每行一个paragraph"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    
    for i, (text, opts) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        fs = opts.get('size', font_size)
        c = opts.get('color', color)
        b = opts.get('bold', bold_first and i == 0)
        
        p.text = text
        p.font.size = Pt(fs)
        p.font.color.rgb = c
        p.font.bold = b
        p.font.name = font_name
        p.space_after = Pt(2)
        p.line_spacing = Pt(fs * line_spacing)
    
    return txBox

def add_header_bar(slide, title_text, subtitle_text=""):
    """统一顶部标题栏"""
    # 深色顶栏
    bar = add_rect(slide, 0, 0, SLIDE_W, 600000, fill_color=BG_HEADER)
    # 底部细线
    line = add_rect(slide, 0, 600000, SLIDE_W, 25000, fill_color=BLUE)
    
    # 标题
    add_text_box(slide, MARGIN, 130000, 7000000, 350000, title_text, 
                 font_size=24, color=WHITE, bold=True)
    
    if subtitle_text:
        add_text_box(slide, SLIDE_W - MARGIN - 4000000, 155000, 4000000, 280000, subtitle_text,
                     font_size=11, color=GREY, bold=False, alignment=PP_ALIGN.RIGHT)

def add_kpi_card(slide, left, top, width, label, value, value_color=None, sub_text="", 
                  bg_color=BG_CARD, label_size=9, value_size=14):
    """KPI数据卡片"""
    card = add_rect(slide, left, top, width, 520000, fill_color=bg_color, radius=60000)
    
    # 标签
    add_text_box(slide, left + 80000, top + 60000, width - 160000, 140000, label,
                 font_size=label_size, color=GREY, bold=False)
    
    # 数值
    vc = value_color or WHITE
    add_text_box(slide, left + 80000, top + 190000, width - 160000, 200000, value,
                 font_size=value_size, color=vc, bold=True)
    
    # 副文本
    if sub_text:
        add_text_box(slide, left + 80000, top + 380000, width - 160000, 120000, sub_text,
                     font_size=8, color=GREY, bold=False)

def add_section_title(slide, left, top, width, number, title, color=BLUE):
    """分区标题 + 色块标记"""
    # 色块
    add_rect(slide, left, top, 45000, 280000, fill_color=color)
    # 编号
    add_text_box(slide, left + 90000, top + 10000, 500000, 260000, number,
                 font_size=20, color=color, bold=True)
    # 标题
    add_text_box(slide, left + 520000, top + 10000, width - 520000, 260000, title,
                 font_size=18, color=WHITE, bold=True)

def add_divider(slide, left, top, width):
    """水平分割线"""
    add_rect(slide, left, top, width, 12000, fill_color=BORDER)

# ═══════════════════════════════════════
# 创建 PPT
# ═══════════════════════════════════════

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# 使用空白布局
blank_layout = prs.slide_layouts[6]  # blank

# ─────────────────────────────────────
# 第1页：封面
# ─────────────────────────────────────
slide1 = prs.slides.add_slide(blank_layout)
add_bg(slide1, BG_DARK)

# 装饰性顶部色条
add_rect(slide1, 0, 0, SLIDE_W, 45000, fill_color=BLUE)

# 左侧装饰竖线
add_rect(slide1, MARGIN, 1200000, 40000, 2400000, fill_color=BLUE)

# 主标题
add_text_box(slide1, MARGIN + 180000, 1200000, 9000000, 750000,
             "2026 推广年中会议汇报", font_size=44, color=WHITE, bold=True)

# 英文副标题
add_text_box(slide1, MARGIN + 180000, 1950000, 8000000, 320000,
             "E-commerce Promotion Mid-Year Review", font_size=16, color=GREY, bold=False)

# 标签行
add_text_box(slide1, MARGIN + 180000, 2350000, 7000000, 350000,
             "H1 平台复盘  /  H2 三大规划  /  组织升级", font_size=20, color=BLUE_LIGHT, bold=True)

# 分割线
add_divider(slide1, MARGIN + 180000, 2780000, 6200000)

# 底部信息区
info_y = 4870000
add_text_box(slide1, MARGIN + 180000, info_y, 3000000, 280000, 
             "汇报人：东方", font_size=13, color=GREY)
add_text_box(slide1, MARGIN + 180000, info_y + 320000, 3000000, 280000,
             "部门：国内营销中心", font_size=13, color=GREY)
add_text_box(slide1, MARGIN + 180000, info_y + 640000, 3000000, 280000,
             "汇报时间：2026.07.23", font_size=13, color=GREY)

# 右下角标语
add_text_box(slide1, SLIDE_W - MARGIN - 4000000, SLIDE_H - MARGIN - 350000, 4000000, 350000,
             "从投放执行到增长经营", font_size=16, color=BLUE_LIGHT, bold=True, alignment=PP_ALIGN.RIGHT)

# ─────────────────────────────────────
# 第2页：目录
# ─────────────────────────────────────
slide2 = prs.slides.add_slide(blank_layout)
add_bg(slide2, BG_DARK)
add_header_bar(slide2, "汇报结构", "CONTENTS")

# 三大模块卡片
modules = [
    ("01", "上半年复盘", 
     "先看平台内各店铺同比、环比表现，明确主力店铺与补充店铺的角色。\n重点拆解天猫旗舰店和京东自营的效率矛盾与优化方向。",
     BLUE, BLUE_BG),
    ("02", "下半年规划",
     "围绕商品分层、品类增长和渠道策略，形成可执行的投放打法。\n把预算、人群、关键词、素材和止损规则落到具体经营场景。",
     GREEN, GREEN_BG),
    ("03", "组织升级",
     "通过成员能力建设、协同机制和 AI 工具化，降低对个人经验的依赖。\n让数据诊断、策略输出和复盘沉淀成为团队稳定能力。",
     PURPLE, PURPLE_BG),
]

card_w = 3400000
card_h = 2300000
card_gap = 250000
total_w = card_w * 3 + card_gap * 2
start_x = (SLIDE_W - total_w) // 2
card_y = 900000

for idx, (num, title, desc, accent, bg) in enumerate(modules):
    x = start_x + idx * (card_w + card_gap)
    
    # 卡片背景
    add_rect(slide2, x, card_y, card_w, card_h, fill_color=bg, radius=120000)
    # 顶部色条
    add_rect(slide2, x, card_y, card_w, 50000, fill_color=accent)
    
    # 编号
    add_text_box(slide2, x + 200000, card_y + 150000, card_w - 400000, 400000, num,
                 font_size=42, color=accent, bold=True)
    # 标题
    add_text_box(slide2, x + 200000, card_y + 500000, card_w - 400000, 350000, title,
                 font_size=22, color=WHITE, bold=True)
    # 描述
    add_text_box(slide2, x + 200000, card_y + 900000, card_w - 400000, 1200000, desc,
                 font_size=13, color=GREY, bold=False)

# 底部主线说明
mainline_y = card_y + card_h + 350000
add_rect(slide2, start_x, mainline_y, total_w, 750000, fill_color=BG_ACCENT, radius=100000)

add_text_box(slide2, start_x + 250000, mainline_y + 120000, 2200000, 300000, "汇报主线",
             font_size=17, color=BLUE_LIGHT, bold=True)
add_text_box(slide2, start_x + 250000, mainline_y + 420000, total_w - 500000, 300000,
             "从上半年平台与店铺表现出发，识别增长效率和费用节奏问题；下半年通过商品分层、品类协同和渠道差异化策略，推动推广从执行动作升级为增长经营能力。",
             font_size=12, color=GREY)

# ─────────────────────────────────────
# 第3页：天猫店铺表现
# ─────────────────────────────────────
slide3 = prs.slides.add_slide(blank_layout)
add_bg(slide3, BG_DARK)
add_header_bar(slide3, "上半年复盘 01：天猫店铺表现", "主店突出销售同比/环比，ROI 用百分比变化")

# 说明文字
add_text_box(slide3, MARGIN, 710000, 10000000, 260000,
             "天猫板块以旗舰店为核心展开：主店第一行突出效率与销售变化，数码旗舰店只做新起盘补充；常规复盘不展示销售、消耗、费比金额。",
             font_size=11, color=GREY)

# ── 天猫旗舰店 KPI 行 ──
shop_y = 1000000
label_w = 1600000
kpi_w = 1080000

# 店名标签卡片
add_rect(slide3, MARGIN, shop_y, label_w, 1100000, fill_color=BLUE_BG, radius=80000)
add_text_box(slide3, MARGIN + 100000, shop_y + 150000, label_w - 200000, 250000, "天猫旗舰店",
             font_size=16, color=WHITE, bold=True)
add_text_box(slide3, MARGIN + 100000, shop_y + 450000, label_w - 200000, 200000, "重点店铺",
             font_size=10, color=BLUE_LIGHT, bold=True)
add_text_box(slide3, MARGIN + 100000, shop_y + 700000, label_w - 200000, 300000, 
             "主线展开：效率变化 + 费用节奏 + 销售变化", font_size=9, color=GREY)

# KPI 数据
kpi_data_tmall = [
    ("消耗同比", "+30.3%", BLUE_LIGHT),
    ("消耗环比", "-0.8%", GREY),
    ("ROI同比", "+39.3%", GREEN),
    ("ROI环比", "+10.4%", GREEN),
    ("费比同比", "+2.6pt", AMBER),
    ("费比环比", "-0.9pt", GREEN),
    ("销售同比", "+10.0%", GREEN),
    ("销售环比", "+4.4%", GREEN),
]

for i, (label, value, vc) in enumerate(kpi_data_tmall):
    x = MARGIN + label_w + 80000 + i * (kpi_w + 40000)
    add_kpi_card(slide3, x, shop_y, kpi_w, label, value, value_color=vc)

# ── 天猫数码旗舰店 KPI 行 ──
digital_y = 2200000
add_rect(slide3, MARGIN, digital_y, label_w, 750000, fill_color=BG_ACCENT, radius=80000)
add_text_box(slide3, MARGIN + 100000, digital_y + 150000, label_w - 200000, 220000, "天猫数码旗舰店",
             font_size=14, color=WHITE, bold=True)
add_text_box(slide3, MARGIN + 100000, digital_y + 400000, label_w - 200000, 250000, "新起盘账户",
             font_size=10, color=GREY)

digital_labels = ["消耗同比", "消耗环比", "ROI同比", "ROI环比", "费比同比", "费比环比", "CPC同比", "CVR同比"]
for i, label in enumerate(digital_labels):
    x = MARGIN + label_w + 80000 + i * (kpi_w + 40000)
    add_kpi_card(slide3, x, digital_y, kpi_w, label, "新起盘", value_color=GREY,
                  bg_color=BG_ACCENT, label_size=8, value_size=11)

# ── 三个结论卡片 ──
conclusion_y = 3150000
conc_w = 3200000
conc_gap = 250000
conc_start = MARGIN + 500000

conclusions_tmall = [
    ("推广优化效果", "ROI环比 +10.4%", "2026H1 vs 2025H2，效率修复明显", GREEN, GREEN_BG),
    ("同销售费比", "环比 -0.9pt", "费用效率改善，控费提效成立", GREEN, GREEN_BG),
    ("Q2 销售例外说明", "同比 +9.6% / 环比 +24.4%", "平台联合高费比投入带动销售增长明显", BLUE_LIGHT, BLUE_BG),
]

for i, (label, value, sub, vc, bg) in enumerate(conclusions_tmall):
    x = conc_start + i * (conc_w + conc_gap)
    add_rect(slide3, x, conclusion_y, conc_w, 800000, fill_color=bg, radius=100000)
    add_text_box(slide3, x + 180000, conclusion_y + 120000, conc_w - 360000, 200000, label,
                 font_size=11, color=GREY)
    add_text_box(slide3, x + 180000, conclusion_y + 320000, conc_w - 360000, 200000, value,
                 font_size=18, color=vc, bold=True)
    add_text_box(slide3, x + 180000, conclusion_y + 540000, conc_w - 360000, 200000, sub,
                 font_size=9, color=GREY)

# ── 诊断区 ──
diag_y = 4200000
diag_w = 5000000

# 天猫旗舰店诊断
add_rect(slide3, MARGIN, diag_y, diag_w, 1100000, fill_color=BG_CARD, radius=100000)
add_section_title(slide3, MARGIN + 150000, diag_y + 80000, diag_w - 300000, "", "天猫旗舰店：重点诊断", BLUE)
add_multiline_text(slide3, MARGIN + 150000, diag_y + 400000, diag_w - 300000, 650000, [
    ("ROI 同比 +39.3%、环比 +10.4%，推广优化效果明确。", {}),
    ("消耗同比 +30.3%，说明平台联合与活动资源带来放量，后续要控制预算节奏。", {}),
    ("主店 H2 重点：保爆款、控节奏、抢核心词，站外引流必须回看站内承接。", {}),
], font_size=11, color=GREY)

# 数码旗舰店诊断
add_rect(slide3, MARGIN + diag_w + 200000, diag_y, diag_w + 300000, 1100000, fill_color=BG_CARD, radius=100000)
add_section_title(slide3, MARGIN + diag_w + 350000, diag_y + 80000, diag_w, "", "天猫数码旗舰店：稍带过", PURPLE)
add_multiline_text(slide3, MARGIN + diag_w + 350000, diag_y + 400000, diag_w, 650000, [
    ("新起盘账户，缺少完整同比/环比评价基础。", {}),
    ("短期重点看模型冷启、人群沉淀、素材跑数和商品承接。", {}),
    ("H2 以轻预算测试为主，先建规则，再看扩量。", {}),
], font_size=11, color=GREY)

# 底部结论
add_rect(slide3, MARGIN, 5600000, SLIDE_W - MARGIN * 2, 500000, fill_color=BG_ACCENT, radius=80000)
add_text_box(slide3, MARGIN + 200000, 5660000, SLIDE_W - MARGIN * 2 - 400000, 380000,
             '天猫结论：旗舰店是主线，讲清「Q2 销售增长明显，同时 H1 ROI 环比提升、费比环比下降」；数码旗舰店只作为新起盘补充。',
             font_size=12, color=WHITE, bold=True)

# ─────────────────────────────────────
# 第4页：京东店铺表现
# ─────────────────────────────────────
slide4 = prs.slides.add_slide(blank_layout)
add_bg(slide4, BG_DARK)
add_header_bar(slide4, "上半年复盘 02：京东店铺表现", "主店突出销售同比/环比，ROI 用百分比变化")

add_text_box(slide4, MARGIN, 710000, 10000000, 260000,
             "京东板块以自营为核心展开：主店第一行突出效率与销售变化，POP、Mr.Green 作为小店补充；常规复盘不展示销售、消耗、费比金额。",
             font_size=11, color=GREY)

# ── 京东自营 KPI 行 ──
shop_y = 1000000

add_rect(slide4, MARGIN, shop_y, label_w, 1100000, fill_color=AMBER_BG, radius=80000)
add_text_box(slide4, MARGIN + 100000, shop_y + 150000, label_w - 200000, 250000, "京东自营",
             font_size=16, color=WHITE, bold=True)
add_text_box(slide4, MARGIN + 100000, shop_y + 450000, label_w - 200000, 200000, "重点店铺",
             font_size=10, color=AMBER, bold=True)
add_text_box(slide4, MARGIN + 100000, shop_y + 700000, label_w - 200000, 300000,
             "主线展开：效率变化 + 费用节奏 + 销售变化", font_size=9, color=GREY)

kpi_data_jd = [
    ("消耗同比", "-13.2%", GREY),
    ("消耗环比", "-15.8%", GREY),
    ("ROI同比", "+25.3%", GREEN),
    ("ROI环比", "+36.4%", GREEN),
    ("费比同比", "+0.3pt", AMBER),
    ("费比环比", "-1.0pt", GREEN),
    ("销售同比", "-15.8%", RED),
    ("销售环比", "-8.1%", RED),
]

for i, (label, value, vc) in enumerate(kpi_data_jd):
    x = MARGIN + label_w + 80000 + i * (kpi_w + 40000)
    add_kpi_card(slide4, x, shop_y, kpi_w, label, value, value_color=vc)

# ── POP & Mr.Green ──
sub_y = 2200000
sub_labels = ["消耗同比", "消耗环比", "ROI同比", "ROI环比", "费比同比", "费比环比", "CPC同比", "CVR同比"]

# POP
add_rect(slide4, MARGIN, sub_y, label_w, 750000, fill_color=BG_ACCENT, radius=80000)
add_text_box(slide4, MARGIN + 100000, sub_y + 150000, label_w - 200000, 220000, "京东POP",
             font_size=14, color=WHITE, bold=True)

pop_values = ["+3.2%", "-10.8%", "+26.3%", "+35.1%", "+0.1pt", "-0.6pt", "-17.9%", "+12.9%"]
for i, (label, value) in enumerate(zip(sub_labels, pop_values)):
    x = MARGIN + label_w + 80000 + i * (kpi_w + 40000)
    vc = GREEN if value.startswith('+') and ('ROI' in label or 'CVR' in label or '环比' in label) else \
         RED if value.startswith('-') and ('CPC' not in label) else GREY
    add_kpi_card(slide4, x, sub_y, kpi_w, label, value, value_color=vc,
                  bg_color=BG_ACCENT, label_size=8, value_size=11)

# Mr.Green
mr_y = 3050000
add_rect(slide4, MARGIN, mr_y, label_w, 750000, fill_color=BG_ACCENT, radius=80000)
add_text_box(slide4, MARGIN + 100000, mr_y + 150000, label_w - 200000, 220000, "Mr.Green自营",
             font_size=14, color=WHITE, bold=True)

mr_values = ["-19.3%", "-7.2%", "+8.5%", "+6.5%", "-0.9pt", "-0.7pt", "-12.2%", "+2.1%"]
for i, (label, value) in enumerate(zip(sub_labels, mr_values)):
    x = MARGIN + label_w + 80000 + i * (kpi_w + 40000)
    vc = GREEN if value.startswith('+') and ('ROI' in label or 'CVR' in label) else GREY
    add_kpi_card(slide4, x, mr_y, kpi_w, label, value, value_color=vc,
                  bg_color=BG_ACCENT, label_size=8, value_size=11)

# ── 结论卡片 ──
conclusion_y = 4000000

conclusions_jd = [
    ("推广优化效果", "ROI环比 +36.4%", "2026H1 vs 2025H2，自营修复最明显", GREEN, GREEN_BG),
    ("同销售费比", "环比 -1.0pt", "费用效率改善，控费提效明确", GREEN, GREEN_BG),
    ("销售弹性问题", "同比 -15.8% / 环比 -8.1%", "效率提升后，H2 需要恢复核心品类扩量", AMBER, AMBER_BG),
]

for i, (label, value, sub, vc, bg) in enumerate(conclusions_jd):
    x = conc_start + i * (conc_w + conc_gap)
    add_rect(slide4, x, conclusion_y, conc_w, 800000, fill_color=bg, radius=100000)
    add_text_box(slide4, x + 180000, conclusion_y + 120000, conc_w - 360000, 200000, label,
                 font_size=11, color=GREY)
    add_text_box(slide4, x + 180000, conclusion_y + 320000, conc_w - 360000, 200000, value,
                 font_size=18, color=vc, bold=True)
    add_text_box(slide4, x + 180000, conclusion_y + 540000, conc_w - 360000, 200000, sub,
                 font_size=9, color=GREY)

# ── 诊断区 ──
diag_y = 5020000

# 京东自营诊断
add_rect(slide4, MARGIN, diag_y, diag_w, 1100000, fill_color=BG_CARD, radius=100000)
add_section_title(slide4, MARGIN + 150000, diag_y + 80000, diag_w - 300000, "", "京东自营：重点诊断", AMBER)
add_multiline_text(slide4, MARGIN + 150000, diag_y + 400000, diag_w - 300000, 650000, [
    ("ROI 同比 +25.3%、环比 +36.4%，H1 推广优化效果非常明确。", {}),
    ("销售同比、环比仍下滑，说明不能只控费，需要在高 ROI 计划中恢复增长弹性。", {}),
    ("H2 重点：关键词质量、人群分层、爆品预算保护、低效计划止损。", {}),
], font_size=11, color=GREY)

# POP/Mr.Green 诊断
add_rect(slide4, MARGIN + diag_w + 200000, diag_y, diag_w + 300000, 1100000, fill_color=BG_CARD, radius=100000)
add_section_title(slide4, MARGIN + diag_w + 350000, diag_y + 80000, diag_w, "", "POP / Mr.Green：稍带过", PURPLE)
add_multiline_text(slide4, MARGIN + diag_w + 350000, diag_y + 400000, diag_w, 650000, [
    ("京东 POP：ROI 同比 +26.3%、环比 +35.1%，按转化稳定性筛投入。", {}),
    ("Mr.Green 自营：ROI 同比 +8.5%、环比 +6.5%，费比同比/环比均下降，可承担利润型补充。", {}),
    ("两者不作为主叙事，重点用于补充京东平台结构。", {}),
], font_size=11, color=GREY)

# 底部结论
add_rect(slide4, MARGIN, 6350000, SLIDE_W - MARGIN * 2, 420000, fill_color=BG_ACCENT, radius=80000)
add_text_box(slide4, MARGIN + 200000, 6410000, SLIDE_W - MARGIN * 2 - 400000, 300000,
             '京东结论：自营是主线，讲清「ROI 提升明显、同销售费比下降明显，但销售弹性仍需恢复」；POP、Mr.Green 作为结构补充。',
             font_size=12, color=WHITE, bold=True)

# ─────────────────────────────────────
# 第5页：推广精细化提效
# ─────────────────────────────────────
slide5 = prs.slides.add_slide(blank_layout)
add_bg(slide5, BG_DARK)
add_header_bar(slide5, "下半年规划 01：推广精细化提效", "从核心品类关注到全维度监控调优")

# 背景说明
add_text_box(slide5, MARGIN, 740000, 10600000, 260000,
             "背景：过去 1 人统筹多店铺、多品类、多运营模式，精力只能优先看核心品类和重点商品；H2 通过推广助理增加与 AI 工具辅助，把监控范围扩展到全品类、全商品层级、全计划维度。",
             font_size=11, color=GREY)

# 三列卡片
three_cards = [
    ("1. 监控范围升级", BLUE, BLUE_BG,
     ["过去：核心品类、爆款商品、主要账户优先，长尾品和细分计划难持续跟踪。",
      "升级：覆盖平台、店铺、品类、商品、计划、关键词、人群、素材。",
      "目标：从「事后发现大问题」转向「过程识别小异常」。"]),
    ("2. 人力分工升级", GREEN, GREEN_BG,
     ["推广负责人：策略判断、预算分配、平台打法、关键问题决策。",
      "推广助理：数据整理、基础监控、异常标记、素材/计划跟进。",
      "目标：把负责人从重复拉表中释放出来，把时间投入策略和复盘。"]),
    ("3. AI 工具提效", PURPLE, PURPLE_BG,
     ["AI 负责多维数据处理、异常识别、趋势归因和策略草案。",
      "人工负责业务判断、动作取舍、执行落地和结果校准。",
      "目标：提升诊断速度，降低对个人经验的依赖。"]),
]

tc_w = 3400000
tc_h = 1800000
tc_gap = 200000
tc_total = tc_w * 3 + tc_gap * 2
tc_start = (SLIDE_W - tc_total) // 2

for i, (title, accent, bg, items) in enumerate(three_cards):
    x = tc_start + i * (tc_w + tc_gap)
    y = 1150000
    
    add_rect(slide5, x, y, tc_w, tc_h, fill_color=bg, radius=120000)
    add_rect(slide5, x, y, tc_w, 45000, fill_color=accent)
    
    add_text_box(slide5, x + 200000, y + 130000, tc_w - 400000, 280000, title,
                 font_size=16, color=WHITE, bold=True)
    
    item_lines = []
    for item in items:
        item_lines.append((f"• {item}", {}))
    
    add_multiline_text(slide5, x + 200000, y + 480000, tc_w - 400000, 1200000,
                       item_lines, font_size=11, color=GREY)

# 四象限商品分层
quad_y = 3200000
quad_w = 2620000
quad_h = 1500000
quad_gap = 200000
quad_start = MARGIN

quads = [
    ("核心爆款", "保排名、保转化、保预算。", "重点看 ROI 环比、核心词排名、活动承接。", BLUE, BLUE_BG),
    ("潜力单品", "测人群、测素材、测场景。", "重点看 CPC、CVR、加购和 ROI 趋势。", GREEN, GREEN_BG),
    ("新品", "小预算冷启，先跑模型。", "重点看点击率、加购率、素材胜出率。", AMBER, AMBER_BG),
    ("中长尾品", "低预算验证，快速止损。", "重点看 ROI 阈值、转化率、费用泄露。", PURPLE, PURPLE_BG),
]

for i, (title, line1, line2, accent, bg) in enumerate(quads):
    x = quad_start + i * (quad_w + quad_gap)
    
    add_rect(slide5, x, quad_y, quad_w, quad_h, fill_color=bg, radius=100000)
    add_text_box(slide5, x + 150000, quad_y + 120000, quad_w - 300000, 280000, title,
                 font_size=15, color=accent, bold=True)
    add_text_box(slide5, x + 150000, quad_y + 450000, quad_w - 300000, 250000, line1,
                 font_size=11, color=WHITE)
    add_text_box(slide5, x + 150000, quad_y + 750000, quad_w - 300000, 500000, line2,
                 font_size=10, color=GREY)

# 底部落地机制
mech_y = 5000000
add_rect(slide5, MARGIN, mech_y, SLIDE_W - MARGIN * 2, 600000, fill_color=BG_ACCENT, radius=100000)
add_text_box(slide5, MARGIN + 200000, mech_y + 80000, SLIDE_W - MARGIN * 2 - 400000, 400000,
             "落地机制：日看异常、周看商品层级迁移、月看平台结构；推广助理负责数据与跟进，AI 负责初筛和策略草案，负责人负责判断和关键动作。",
             font_size=12, color=WHITE, bold=True)

# ─────────────────────────────────────
# 第6页：赋能品类增长
# ─────────────────────────────────────
slide6 = prs.slides.add_slide(blank_layout)
add_bg(slide6, BG_DARK)
add_header_bar(slide6, "下半年规划 02：赋能品类增长", "通过推广数据反推经营动作")

add_text_box(slide6, MARGIN, 740000, 10600000, 300000,
             "目标：推广不只是买流量，而是通过投放组合和数据反馈，帮助品类找到更确定的增长路径，指导货品、价格、内容和活动节奏。",
             font_size=12, color=GREY)

# 三列大卡片
big_three = [
    ("1. 数据反馈", "识别机会", GREEN,
     ["按品类/商品监控 ROI、CPC、CVR、点击率、收藏加购和转化趋势。",
      "识别三类机会：低 CPC 高 CVR 的放量机会；高点击低转化的承接问题；高消耗低 ROI 的止损对象。",
      "输出给运营：哪些品类值得加预算，哪些商品需要优化页面、价格或权益。"]),
    ("2. 组合打法", "提升效率", BLUE,
     ["搜索负责承接明确需求，推荐/场景负责拓新人群，站外种草负责补认知。",
      "核心品做搜索+推荐联动，潜力品做小预算多场景测试，新品先内容种草再站内承接。",
      "每类打法都要绑定复盘指标，避免只看投放动作、不看经营结果。"]),
    ("3. 经营协同", "放大结果", PURPLE,
     ["推广结论同步给商品、运营、供应链：爆款备货、价格策略、详情页承接、活动资源申请。",
      "大促前做货品池和预算池联审，避免推广有流量但库存、权益、页面承接不足。",
      "沉淀品类推广打法，形成可复用的增长模型。"]),
]

bc_w = 3400000
bc_h = 3700000
bc_gap = 200000
bc_total = bc_w * 3 + bc_gap * 2
bc_start = (SLIDE_W - bc_total) // 2

for i, (title, subtitle, accent, items) in enumerate(big_three):
    x = bc_start + i * (bc_w + bc_gap)
    y = 1150000
    
    add_rect(slide6, x, y, bc_w, bc_h, fill_color=BG_CARD, radius=120000)
    add_rect(slide6, x, y, bc_w, 50000, fill_color=accent)
    
    add_text_box(slide6, x + 250000, y + 150000, bc_w - 500000, 300000, title,
                 font_size=20, color=accent, bold=True)
    add_text_box(slide6, x + 250000, y + 420000, bc_w - 500000, 250000, subtitle,
                 font_size=14, color=WHITE, bold=True)
    
    item_lines = []
    for idx, item in enumerate(items):
        item_lines.append((f"{idx+1}. {item}", {'size': 11, 'color': GREY}))
    
    add_multiline_text(slide6, x + 250000, y + 750000, bc_w - 500000, 2800000,
                       item_lines, font_size=11, color=GREY, line_spacing=1.5)

# 底部结论
bot_y = 5150000
add_rect(slide6, MARGIN, bot_y, SLIDE_W - MARGIN * 2, 450000, fill_color=BG_ACCENT, radius=80000)
add_text_box(slide6, MARGIN + 200000, bot_y + 60000, SLIDE_W - MARGIN * 2 - 400000, 300000,
             '核心输出：从「投放复盘表」升级为「品类增长建议表」，让推广数据真正指导经营动作。',
             font_size=13, color=WHITE, bold=True)

# ─────────────────────────────────────
# 第7页：渠道策略方向
# ─────────────────────────────────────
slide7 = prs.slides.add_slide(blank_layout)
add_bg(slide7, BG_DARK)
add_header_bar(slide7, "下半年规划 03：渠道策略方向", "天猫 / 京东 / 拼多多")

add_text_box(slide7, MARGIN, 740000, 10600000, 300000,
             "核心逻辑：渠道策略不再用同一套 ROI 标准一刀切，而是按平台角色分配目标、预算和考核指标；天猫抢市场，京东提效率，拼多多先诊断后扩量。",
             font_size=11, color=GREY)

channels = [
    ("天猫", "平台联合 + 高费比阶段抢占", BLUE, BLUE_BG,
     ["角色定位：品牌主阵地和市场份额承接渠道。",
      "策略重点：围绕平台联合资源、活动节点和核心爆款，阶段性接受高费比换取市场曝光与销售增长。",
      "执行动作：核心词卡位、推荐扩量、UD 站外引流、内容种草回流。",
      "管理边界：高费比必须绑定活动目标、商品池和承接结果，避免无目标放量。"]),
    ("京东", "效率优先 + 爆品恢复弹性", AMBER, AMBER_BG,
     ["角色定位：效率修复和利润稳定渠道。",
      "策略重点：以京东自营为主，守住高 ROI 计划，同时恢复核心品类与爆品扩量。",
      "执行动作：关键词质量分、人群分层、低效词清理、爆品预算保护。",
      "管理边界：优先看 ROI 环比、费比环比、CPC/CVR 变化，防止只控费不增长。"]),
    ("拼多多", "账户诊断 + 规则先行", PURPLE, PURPLE_BG,
     ["角色定位：9 月接手后的新增提效场。",
      "策略重点：不急于规模放量，先完成账户体检、商品分层、计划结构和投放规则搭建。",
      "执行动作：搜索/场景基础计划、可投商品池、素材测试、ROI 阈值和低效止损机制。",
      "管理边界：先跑通方法，再放大预算，避免接手初期粗放消耗。"]),
]

ch_w = 3400000
ch_h = 4300000
ch_start = (SLIDE_W - ch_w * 3 - tc_gap * 2) // 2

for i, (name, subtitle, accent, bg, items) in enumerate(channels):
    x = ch_start + i * (ch_w + tc_gap)
    y = 1150000
    
    add_rect(slide7, x, y, ch_w, ch_h, fill_color=bg, radius=120000)
    add_rect(slide7, x, y, ch_w, 55000, fill_color=accent)
    
    # 平台名
    add_text_box(slide7, x + 200000, y + 120000, ch_w - 400000, 300000, name,
                 font_size=22, color=accent, bold=True)
    # 副标题
    add_text_box(slide7, x + 200000, y + 400000, ch_w - 400000, 250000, subtitle,
                 font_size=14, color=WHITE, bold=True)
    
    item_lines = []
    for item in items:
        item_lines.append((f"▸ {item}", {'size': 11, 'color': GREY}))
    
    add_multiline_text(slide7, x + 200000, y + 750000, ch_w - 400000, 3400000,
                       item_lines, font_size=11, color=GREY, line_spacing=1.45)

# 底部判断
bot_y = 5700000
add_rect(slide7, MARGIN, bot_y, SLIDE_W - MARGIN * 2, 500000, fill_color=BG_ACCENT, radius=80000)
add_text_box(slide7, MARGIN + 200000, bot_y + 70000, SLIDE_W - MARGIN * 2 - 400000, 350000,
             "渠道判断：天猫用资源换增长，但要看承接；京东用效率守利润，但要补弹性；拼多多先建规则，再进入规模化提效。",
             font_size=13, color=WHITE, bold=True)

# ─────────────────────────────────────
# 第8页：组织成员管理
# ─────────────────────────────────────
slide8 = prs.slides.add_slide(blank_layout)
add_bg(slide8, BG_DARK)
add_header_bar(slide8, "组织升级 01：组织成员管理", "能力建设 / 分工机制 / 降低关键人员依赖")

add_text_box(slide8, MARGIN, 740000, 10600000, 300000,
             "核心目标：把推广负责人个人经验拆成团队可执行的方法，让推广助理、运营和品类负责人能围绕同一套指标协作，逐步降低业务对单一关键人员经验的依赖。",
             font_size=11, color=GREY)

org_cards = [
    ("1. 能力短板", "从操作能力到经营判断", BLUE,
     ["当前问题：成员能做基础执行，但对平台机制、商品层级和预算节奏的判断不足。",
      "能力方向：从「会建计划、会调价」升级为「会看趋势、会归因、会判断扩量/止损」。",
      "培养重点：ROI/费比/CPC/CVR 联动分析，商品池分层，活动节奏与推广节奏匹配。"]),
    ("2. 分工机制", "负责人 + 助理协同", GREEN,
     ["推广负责人：负责平台策略、预算分配、关键账户诊断、跨部门协同。",
      "推广助理：负责日报数据整理、异常初筛、基础计划跟进、素材和关键词清单维护。",
      "协同方式：助理先筛异常，负责人判断优先级并给动作，复盘后沉淀为 SOP。"]),
    ("3. 标准沉淀", "减少经验黑箱", PURPLE,
     ["建立商品分层 SOP：核心爆款、潜力单品、新品、中长尾品分别对应投放规则。",
      "建立异常处理清单：ROI 下滑、费比上升、CPC 异常、CVR 波动、消耗突增对应排查路径。",
      "建立复盘机制：周看计划、月看平台、活动看商品和承接，把经验转成团队资产。"]),
]

for i, (title, subtitle, accent, items) in enumerate(org_cards):
    x = tc_start + i * (tc_w + tc_gap)
    y = 1150000
    
    add_rect(slide8, x, y, tc_w, 3800000, fill_color=BG_CARD, radius=120000)
    add_rect(slide8, x, y, tc_w, 45000, fill_color=accent)
    
    add_text_box(slide8, x + 200000, y + 130000, tc_w - 400000, 300000, title,
                 font_size=18, color=accent, bold=True)
    add_text_box(slide8, x + 200000, y + 400000, tc_w - 400000, 250000, subtitle,
                 font_size=14, color=WHITE, bold=True)
    
    item_lines = []
    for item in items:
        item_lines.append((f"• {item}", {'size': 11, 'color': GREY}))
    
    add_multiline_text(slide8, x + 200000, y + 750000, tc_w - 400000, 2850000,
                       item_lines, font_size=11, color=GREY, line_spacing=1.4)

# 底部结论
add_rect(slide8, MARGIN, 5250000, SLIDE_W - MARGIN * 2, 450000, fill_color=BG_ACCENT, radius=80000)
add_text_box(slide8, MARGIN + 200000, 5310000, SLIDE_W - MARGIN * 2 - 400000, 300000,
             "组织升级不是增加人手本身，而是用清晰分工和标准流程，把团队从「被动执行」推向「主动诊断、主动优化」。",
             font_size=13, color=WHITE, bold=True)

# ─────────────────────────────────────
# 第9页：AI 赋能团队
# ─────────────────────────────────────
slide9 = prs.slides.add_slide(blank_layout)
add_bg(slide9, BG_DARK)
add_header_bar(slide9, "组织升级 02：AI 赋能团队", "AI 策略与实际工作结合")

add_text_box(slide9, MARGIN, 740000, 10600000, 300000,
             "目标：AI 不替代推广判断，而是提升数据处理、问题定位和策略草案效率；人工负责业务判断、动作取舍、执行落地和复盘纠偏。",
             font_size=12, color=GREY)

ai_cards = [
    ("1. AI 多维数据处理", PURPLE,
     ["把平台、店铺、商品、计划、关键词数据结构化，自动识别消耗、ROI、CPC、CVR 的异常变化。",
      "从「人工拉表找问题」升级为「AI 先筛异常，人工判断原因」。",
      "优先建立日报异常提醒和周度复盘摘要。"]),
    ("2. AI 快速定位核心问题", BLUE,
     ["由总到细：先看平台，再看店铺，再看商品，再看计划/关键词。",
      "建立核心指标模块：消耗、ROI、CPC、CVR、点击率、加购、转化。",
      "让 AI 输出可能原因：流量成本、转化承接、预算节奏、商品竞争力、活动资源。"]),
    ("3. AI 专业策略输出", GREEN,
     ["沉淀推广指令库：不同异常场景对应不同分析指令和优化建议。",
      "运营提出组合问题，AI 生成策略草案，推广负责人校验后指导实操。",
      "降低推广指导成本，让团队更快形成标准化判断。"]),
]

for i, (title, accent, items) in enumerate(ai_cards):
    x = tc_start + i * (tc_w + tc_gap)
    y = 1200000
    
    add_rect(slide9, x, y, tc_w, 3800000, fill_color=BG_CARD, radius=120000)
    add_rect(slide9, x, y, tc_w, 45000, fill_color=accent)
    
    add_text_box(slide9, x + 200000, y + 150000, tc_w - 400000, 300000, title,
                 font_size=18, color=accent, bold=True)
    
    item_lines = []
    for item in items:
        item_lines.append((f"▸ {item}", {'size': 11, 'color': GREY}))
    
    add_multiline_text(slide9, x + 200000, y + 550000, tc_w - 400000, 3000000,
                       item_lines, font_size=11, color=GREY, line_spacing=1.4)

# AI 闭环
close_y = 5300000
add_rect(slide9, MARGIN, close_y, SLIDE_W - MARGIN * 2, 400000, fill_color=BG_ACCENT, radius=80000)
add_text_box(slide9, MARGIN + 200000, close_y + 50000, SLIDE_W - MARGIN * 2 - 400000, 280000,
             "AI 闭环：数据清洗 → 异常识别 → 策略草案 → 人工判断 → 执行复盘 → 指令库沉淀",
             font_size=13, color=WHITE, bold=True)

# ─────────────────────────────────────
# 第10页：行动节奏与协同
# ─────────────────────────────────────
slide10 = prs.slides.add_slide(blank_layout)
add_bg(slide10, BG_DARK)
add_header_bar(slide10, "H2 行动节奏与协同机制", "MILESTONE & COLLABORATION")

add_text_box(slide10, MARGIN, 740000, 10600000, 260000,
             "把 H2 规划落到时间表：先建规则，再接渠道，再打大促，最后沉淀年度方法。",
             font_size=12, color=GREY)

# 时间轴四阶段
timeline = [
    ("7-8 月", BLUE,
     ["完善商品分层规则、周度复盘模板和异常处理清单。",
      "明确天猫、京东预算阈值和重点商品池。"]),
    ("9 月", PURPLE,
     ["接手拼多多，完成账户诊断、商品分层、计划结构梳理。",
      "建立搜索/场景基础投放规则。"]),
    ("10-11 月", AMBER,
     ["围绕双11执行预算节奏、爆品保护、低效止损。",
      "大促后复盘平台、商品、素材和承接问题。"]),
    ("12 月", GREEN,
     ["沉淀年度推广手册、AI 指令库和品类打法库。",
      "输出 2027 年预算与渠道规划建议。"]),
]

tl_w = 2550000
tl_h = 1500000
tl_gap = 200000
tl_total = tl_w * 4 + tl_gap * 3
tl_start = (SLIDE_W - tl_total) // 2

for i, (month, accent, items) in enumerate(timeline):
    x = tl_start + i * (tl_w + tl_gap)
    y = 1150000
    
    add_rect(slide10, x, y, tl_w, tl_h, fill_color=BG_CARD, radius=120000)
    add_rect(slide10, x, y, tl_w, 50000, fill_color=accent)
    
    add_text_box(slide10, x + 180000, y + 150000, tl_w - 360000, 300000, month,
                 font_size=20, color=accent, bold=True)
    
    item_lines = []
    for item in items:
        item_lines.append((f"• {item}", {'size': 10, 'color': GREY}))
    
    add_multiline_text(slide10, x + 180000, y + 500000, tl_w - 360000, 900000,
                       item_lines, font_size=10, color=GREY, line_spacing=1.35)

# 箭头连接（视觉上）
for i in range(3):
    arrow_x = tl_start + (i + 1) * tl_w + i * tl_gap - 180000
    add_text_box(slide10, arrow_x, 1700000, 200000, 300000, "→",
                 font_size=20, color=GREY_DARK, bold=True, alignment=PP_ALIGN.CENTER)

# 协同机制
collab_y = 3000000
add_rect(slide10, MARGIN, collab_y, SLIDE_W - MARGIN * 2, 1500000, fill_color=BG_CARD, radius=120000)
add_text_box(slide10, MARGIN + 300000, collab_y + 120000, 4000000, 300000, "协同机制",
             font_size=19, color=BLUE_LIGHT, bold=True)

collab_items = [
    ("周度", "推广与运营同步异常和动作"),
    ("月度", "推广输出品类增长建议"),
    ("大促前", "推广、运营、商品、供应链做货品池、预算池、资源位、页面承接联审"),
]

for i, (freq, desc) in enumerate(collab_items):
    item_y = collab_y + 550000 + i * 280000
    # 色标
    colors = [BLUE, GREEN, AMBER]
    add_rect(slide10, MARGIN + 350000, item_y, 35000, 220000, fill_color=colors[i])
    add_text_box(slide10, MARGIN + 500000, item_y, 1500000, 220000, freq,
                 font_size=13, color=colors[i], bold=True)
    add_text_box(slide10, MARGIN + 2000000, item_y, 8000000, 220000, desc,
                 font_size=12, color=GREY)

# ─────────────────────────────────────
# 第11页：致谢
# ─────────────────────────────────────
slide11 = prs.slides.add_slide(blank_layout)
add_bg(slide11, BG_DARK)

# 中央大标题
add_text_box(slide11, 0, 2200000, SLIDE_W, 800000, "THANK YOU",
             font_size=56, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# 分割线
add_rect(slide11, (SLIDE_W - 4000000) // 2, 3150000, 4000000, 25000, fill_color=BLUE)

# 副标题
add_text_box(slide11, 0, 3400000, SLIDE_W, 400000, "让增长更可控，让效率更稳定",
             font_size=18, color=GREY, bold=False, alignment=PP_ALIGN.CENTER)

# 底部信息
add_text_box(slide11, 0, 4800000, SLIDE_W, 300000, "国内营销中心  |  东方  |  2026.07.23",
             font_size=12, color=GREY_DARK, bold=False, alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════
# 保存
# ═══════════════════════════════════════
output_path = "/Users/llano/Desktop/国内hermes存储/东方/2026推广年中会议汇报_优化版.pptx"
prs.save(output_path)
print(f"✅ 已保存: {output_path}")
print(f"共 {len(prs.slides)} 页")
